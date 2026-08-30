#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT生成工具 - 在原模板基础上做数据替换（不重建、不加任何效果）

设计原则：
  1. 原样保留模板所有样式（母版、版式、字体、颜色），只改文字和图片
  2. 只在"硬件选型"页做替换，其它页面一字不动
  3. 支持 --hardware_page 指定替换第几个硬件页（多视角模板每页一个视角）
  4. 绝不写入任何阴影/发光/倒影等效果（历史版本曾被外部合并工具加阴影）

适配的模板占位（以"视觉检测方案.pptx"为例）：
  - 表格行：相机工作距离 / 相机视野 / 相机精度 / 光源工作距离
  - 文本行：使用XXXX万像素相机，分辨率为W × H，单个像素精度为X.XXXXmm/pixel
  - 文本行：远心镜头： XXX / 镜头： XXX
  - 文本行：同轴光： XXX / 环光： XXX
  - 独立文本框：228±4mm（相机WD标注）、80mm（光源WD标注）
  - 右上角相机产品图
"""

import json
import os
import re
import sys
from typing import Dict, List, Optional

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    from pptx.util import Emu
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    print("错误: python-pptx 未安装")
    sys.exit(1)


def set_text_preserving_format(para, new_text: str):
    """设置段落文本，保留第一个run的字体格式（不新增任何效果）"""
    if para.runs:
        first_run = para.runs[0]
        font_name = first_run.font.name
        font_size = first_run.font.size
        font_bold = first_run.font.bold
        font_italic = first_run.font.italic
        font_color = None
        try:
            if first_run.font.color and first_run.font.color.type is not None:
                font_color = first_run.font.color.rgb
        except Exception:
            pass

        for run in para.runs:
            run.text = ""
        first_run.text = new_text

        if font_name:
            first_run.font.name = font_name
        if font_size:
            first_run.font.size = font_size
        if font_bold is not None:
            first_run.font.bold = font_bold
        if font_italic is not None:
            first_run.font.italic = font_italic
        if font_color:
            try:
                first_run.font.color.rgb = font_color
            except Exception:
                pass
    else:
        para.add_run().text = new_text


def is_hardware_slide(slide) -> bool:
    """判断是否硬件选型页：存在短标题文本，等于/包含'硬件选型'等关键词"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text and len(text) < 20 and (
                    text == '硬件选型' or text == '相机选型' or text == '相机参数'):
                return True
    return False


def replace_image_in_slide(slide, image_path: str, position: str = "right_top") -> bool:
    """替换幻灯片指定位置的图片，保持原占位与宽高比"""
    if not os.path.exists(image_path):
        print(f"  警告: 图片文件不存在 - {image_path}")
        return False

    images = [s for s in slide.shapes if s.shape_type == 13]  # PICTURE
    if not images:
        print("  警告: 幻灯片中没有找到图片")
        return False

    target_image = None
    if position == "right_top":
        # x坐标最大且y较小的图片
        best_key = None
        for img in images:
            key = (img.left, -img.top)
            if best_key is None or key > best_key:
                best_key = key
                target_image = img
    elif position == "left_top":
        best_key = None
        for img in images:
            key = (-img.left, -img.top)
            if best_key is None or key > best_key:
                best_key = key
                target_image = img
    else:
        target_image = images[0]

    if target_image is None:
        print("  警告: 未找到目标位置的图片")
        return False

    original_left = target_image.left
    original_top = target_image.top
    original_width = target_image.width
    original_height = target_image.height

    sp = target_image._element
    sp.getparent().remove(sp)

    try:
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        aspect = img_w / img_h
        if original_width / original_height > aspect:
            new_height = original_height
            new_width = int(original_height * aspect)
        else:
            new_width = original_width
            new_height = int(original_width / aspect)

        slide.shapes.add_picture(image_path, original_left, original_top,
                                 new_width, new_height)
        print(f"  已替换图片: {os.path.basename(image_path)}")
        return True
    except Exception as e:
        print(f"  错误: 替换图片失败 - {e}")
        return False


def _format_wan_pixels(res_w, res_h) -> str:
    """分辨率 → 万像素字符串，符合行业叫法取整（501→500万，1996→2000万，1229→1200万）"""
    try:
        mp = int(res_w) * int(res_h) / 10000.0
        step = 100 if mp >= 1000 else 10
        return f"{round(mp / step) * step:.0f}"
    except Exception:
        return ""


def replace_on_hardware_slide(slide, data: Dict) -> int:
    """
    在单个硬件选型页执行全部替换规则，返回替换次数。
    data 键：
      camera_brand, camera_model, camera_res_w, camera_res_h, camera_pixel,
      camera_image_path, lens_brand, lens_model, lens_magnification, lens_wd,
      light_brand, light_model, light_wd, pixel_precision, fov_w, fov_h
    """
    count = 0

    camera_brand = data.get('camera_brand', '')
    camera_model = data.get('camera_model', '')
    camera_res_w = data.get('camera_res_w', '')
    camera_res_h = data.get('camera_res_h', '')
    lens_brand = data.get('lens_brand', '')
    lens_model = data.get('lens_model', '')
    lens_mag = data.get('lens_magnification', '')
    lens_wd = data.get('lens_wd', '')
    # 官网原文公差（如"158±3"）优先于数值——WD标注必须与官网一致，禁止从型号推断
    lens_wd_display = f"{data.get('lens_wd_spec')}mm" if data.get('lens_wd_spec') \
        else (f"{lens_wd}mm" if lens_wd else '')
    light_brand = data.get('light_brand', '')
    light_model = data.get('light_model', '')
    light_wd = data.get('light_wd', '')
    pixel_precision = data.get('pixel_precision', '')
    fov_w = data.get('fov_w', '')
    fov_h = data.get('fov_h', '')

    wan_pixels = _format_wan_pixels(camera_res_w, camera_res_h)

    # ---------- 1. 表格：按行首列标题定向替换 ----------
    for shape in slide.shapes:
        if not shape.has_table:
            continue
        table = shape.table
        for row in table.rows:
            header = row.cells[0].text.strip()
            value_cell = row.cells[1] if len(row.cells) > 1 else None
            if value_cell is None:
                continue
            new_value = None
            if '相机工作距离' in header and lens_wd_display:
                new_value = lens_wd_display
            elif '相机视野' in header and fov_w and fov_h:
                new_value = f"{fov_w}mm×{fov_h}mm"
            elif '相机精度' in header and pixel_precision:
                new_value = f"{pixel_precision}mm/pixel"
            elif '光源工作距离' in header and light_wd:
                new_value = f"{light_wd}mm"

            if new_value and value_cell.text.strip() != new_value:
                for para in value_cell.text_frame.paragraphs:
                    if para.text.strip():
                        set_text_preserving_format(para, new_value)
                        count += 1
                        break

    # ---------- 2. 独立文本框标注（如 228±4mm / 80mm）----------
    # 每页只收集一次：数值大的是相机WD，数值小的是光源WD
    wd_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if re.match(r'^\d+([.。±]\d+)?mm$', text):
                m = re.match(r'^(\d+)', text)
                if m:
                    wd_shapes.append((float(m.group(1)), para, shape))
    if wd_shapes:
        wd_shapes.sort(key=lambda x: x[0], reverse=True)
        # 标注框宽度有限：关自动换行，防止替换后数值变长断字（51.1m/m）
        for _, _, wd_shape in wd_shapes:
            try:
                wd_shape.text_frame.word_wrap = False
            except Exception:
                pass
        if len(wd_shapes) >= 2 and lens_wd_display and light_wd:
            set_text_preserving_format(wd_shapes[0][1], lens_wd_display)
            set_text_preserving_format(wd_shapes[1][1], f"{light_wd}mm")
            count += 2
        elif lens_wd_display:
            set_text_preserving_format(wd_shapes[0][1], lens_wd_display)
            count += 1

    # ---------- 3. 文本段落替换 ----------
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            # 3.1 相机描述行："使用2000万像素相机，分辨率为5472 × 3648，..."
            if re.search(r'使用.*万像素相机', text) and camera_res_w:
                desc = (f"使用{wan_pixels}万像素相机"
                        + (f"（{camera_brand} {camera_model}）" if camera_model else "")
                        + f"，分辨率为{camera_res_w} × {camera_res_h}"
                        + (f"，单个像素精度为{pixel_precision}mm/pixel"
                           if pixel_precision else ""))
                set_text_preserving_format(para, desc)
                count += 1
                continue

            # 3.2 镜头描述行："远心镜头： DTCM110-80H-AL" / "镜头： MVL-..."
            m = re.match(r'^(远心镜头|镜头|双远心镜头)[:：]\s*(.+)$', text)
            if m and lens_model:
                prefix = m.group(1)
                spec = f"{lens_brand} {lens_model}"
                if lens_mag:
                    spec += f"，{lens_mag}x"
                if lens_wd_display:
                    spec += f"，WD {lens_wd_display}"
                set_text_preserving_format(para, f"{prefix}： {spec}")
                count += 1
                continue

            # 3.3 光源描述行："同轴光： FL-COX2-85W" / "环光： FL-R-7090W"
            # 类型前缀必须与实际选中的光源类型一致，不能用模板原前缀
            # （模板第4页原为"同轴光"，选中环形光源时会与第5页"环光"矛盾）
            m = re.match(r'^(同轴光|环光|环形光|背光|光源)[:：]\s*(.+)$', text)
            if m and light_model:
                light_prefix = data.get('light_type') or m.group(1)
                set_text_preserving_format(
                    para, f"{light_prefix}： {light_brand} {light_model}")
                count += 1
                continue

            # 3.3b 示意图中独立的纯类型词标注框（如"同轴光"三个字）：按实际光源类型纠正
            # 标注框窄（约19~25mm），只放短类型词；light_type可能是"LED环形光光源"长描述，
            # 需规范化为"环形光"再写入。同义类型（环光=环形光）且框宽足够时保留模板原词，
            # 避免同长度替换造成溢出（第5页框按2字"环光"设计，写3字会折行溢出）
            if re.match(r'^(同轴光|环光|环形光|背光)$', text) and data.get('light_type'):
                short_type = re.sub(r'^LED|光源$', '', data['light_type']).strip()
                synonyms = ({'环形光', '环光'}, {'同轴光'}, {'背光'})
                same = any(text in g and short_type in g for g in synonyms)
                if short_type and not same and text != short_type:
                    set_text_preserving_format(para, short_type)
                    count += 1
                continue

            # 3.4 兼容旧模板触发词（Basler/acA2500/FL-R-）
            if 'Basler Lens' in text and lens_model:
                set_text_preserving_format(
                    para, re.sub(r'Basler Lens\s+\S+', f"{lens_brand} {lens_model}", text))
                count += 1
                continue
            if 'acA2500' in text and camera_model:
                set_text_preserving_format(
                    para, f"{camera_brand} {camera_model} "
                          f"{data.get('camera_interface', '')} 工业相机")
                count += 1
                continue
            if 'FL-R-' in text and light_model:
                set_text_preserving_format(
                    para, re.sub(r'FL-R-\d+[A-Z]*', f"{light_brand} {light_model}", text))
                count += 1
                continue

            # 3.5 兜底：分辨率数字对（5472 × 3648）
            if camera_res_w and camera_res_h and \
               re.search(r'\d{3,5}\s*[×x]\s*\d{3,5}', text) and '万像素' not in text:
                new_text = re.sub(r'\d{3,5}\s*[×x]\s*\d{3,5}',
                                  f'{camera_res_w} × {camera_res_h}',
                                  text, flags=re.IGNORECASE)
                if new_text != text:
                    set_text_preserving_format(para, new_text)
                    count += 1
                    continue

            # 3.6 兜底：中文万像素描述（500万像素 → 新值）
            if wan_pixels:
                new_text = re.sub(r'\d{2,4}\s*万\s*像素', f'{wan_pixels}万像素', text)
                if new_text != text:
                    set_text_preserving_format(para, new_text)
                    count += 1
                    continue

            # 3.7 兜底：像素精度值（0.0115mm/pixel）
            if pixel_precision and re.search(r'\d\.\d+\s*mm/pix', text):
                new_text = re.sub(r'\d\.\d+\s*mm/pixel', f'{pixel_precision}mm/pixel', text)
                new_text = re.sub(r'\d\.\d+\s*mm/pix', f'{pixel_precision}mm/pixel', new_text)
                if new_text != text:
                    set_text_preserving_format(para, new_text)
                    count += 1
                    continue

            # 3.8 兜底：视野尺寸 63mm×42mm（表格外的正文引用）
            if fov_w and fov_h and re.search(r'\d+mm\s*[×x*]\s*\d+mm', text):
                new_text = re.sub(r'\d+mm\s*[×x*]\s*\d+mm',
                                  f'{fov_w}mm×{fov_h}mm', text)
                if new_text != text:
                    set_text_preserving_format(para, new_text)
                    count += 1
                    continue

    # ---------- 4. 相机产品图 ----------
    image_path = data.get('camera_image_path', '')
    if image_path:
        if replace_image_in_slide(slide, image_path, position="right_top"):
            count += 1

    return count


def generate_ppt(template_path: str,
                 output_path: str,
                 camera: Dict,
                 lens: Dict,
                 light_source: Dict = None,
                 project_name: str = None,
                 field_of_view: Dict = None,
                 hardware_page: int = None,
                 required_fov: Dict = None):
    """
    基于原模板生成PPT：只做数据替换，保留全部样式。

    Args:
        template_path: 模板PPT路径
        output_path: 输出路径
        camera/lens/light_source: 选型结果
        project_name: 项目名（仅用于日志）
        field_of_view: 选型用视野 {'width','height'}
        hardware_page: 替换第几个硬件选型页（1起，None=全部硬件页）
        required_fov: 需求视野（与field_of_view相同时可省略）
    """
    print(f"加载模板: {template_path}")
    prs = Presentation(template_path)

    # ---- 准备替换数据 ----
    fov = field_of_view or required_fov or {}
    actual_fov = lens.get('actual_fov') or {}
    fov_w = fov_h = ''
    # 优先用"实际视野"（相机+镜头闭环计算），保证与精度自洽
    if actual_fov.get('width') and actual_fov.get('height'):
        fov_w = f"{round(actual_fov['width'])}"
        fov_h = f"{round(actual_fov['height'])}"
    elif fov:
        fov_w = f"{round(fov['width'])}"
        fov_h = f"{round(fov['height'])}"

    pixel_precision = ''
    try:
        if camera.get('pixel_size') and lens.get('magnification'):
            pixel_precision = f"{camera['pixel_size'] / 1000.0 / lens['magnification']:.4f}"
    except Exception:
        pass

    # 远心系统工作距离由镜头物方WD决定（相机装远心镜头后不可独立调焦），
    # 因此"相机工作距离"直接取镜头WD；光源WD由选型链按物理约束算好传入
    lens_wd = str(lens.get('working_distance', '') or '')
    light_wd = str(light_source.get('working_distance', '') or '') if light_source else ''
    light_type = str(light_source.get('type', '') or '') if light_source else ''

    # 相机图片绝对路径
    camera_image_path = camera.get('image_path', '')
    if camera_image_path and not os.path.isabs(camera_image_path):
        skill_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        camera_image_path = os.path.join(skill_dir, camera_image_path)

    data = {
        'camera_brand': camera.get('brand', ''),
        'camera_model': camera.get('model', ''),
        'camera_res_w': camera.get('resolution', {}).get('width', ''),
        'camera_res_h': camera.get('resolution', {}).get('height', ''),
        'camera_pixel': camera.get('pixel_size', ''),
        'camera_interface': camera.get('interface', ''),
        'camera_image_path': camera_image_path,
        'lens_brand': lens.get('brand', ''),
        'lens_model': lens.get('model', ''),
        'lens_magnification': lens.get('magnification', ''),
        'lens_wd': lens_wd,
        'lens_wd_spec': str(lens.get('wd_spec', '') or ''),
        'light_brand': light_source.get('brand', '') if light_source else '',
        'light_model': light_source.get('model', '') if light_source else '',
        'light_wd': light_wd,
        'light_type': light_type,
        'pixel_precision': pixel_precision,
        'fov_w': fov_w,
        'fov_h': fov_h,
    }

    # ---- 定位硬件选型页 ----
    hardware_slides = [(idx, slide) for idx, slide in enumerate(prs.slides)
                       if is_hardware_slide(slide)]
    print(f"识别到 {len(hardware_slides)} 个硬件选型页: "
          f"{[i + 1 for i, _ in hardware_slides]}")

    if hardware_page is not None:
        if 1 <= hardware_page <= len(hardware_slides):
            targets = [hardware_slides[hardware_page - 1]]
            print(f"指定替换第 {hardware_page} 个硬件页（第{targets[0][0] + 1}页）")
        else:
            print(f"错误: hardware_page={hardware_page} 超出范围 "
                  f"(1-{len(hardware_slides)})")
            return None
    else:
        targets = hardware_slides

    # ---- 执行替换 ----
    total = 0
    for slide_no, slide in targets:
        n = replace_on_hardware_slide(slide, data)
        print(f"  第{slide_no + 1}页: 替换 {n} 处")
        total += n

    print(f"共替换 {total} 处内容")

    os.makedirs(os.path.dirname(output_path) or '.', exist_ok=True)
    prs.save(output_path)
    print(f"PPT已保存: {output_path}")
    return output_path


def main():
    import argparse

    parser = argparse.ArgumentParser(description='PPT生成工具 - 模板数据替换（保留原样式）')
    parser.add_argument('--template', required=True, help='模板PPT路径')
    parser.add_argument('--output', required=True, help='输出路径')
    parser.add_argument('--camera_json', help='相机参数JSON文件')
    parser.add_argument('--lens_json', help='镜头参数JSON文件')
    parser.add_argument('--light_json', help='光源参数JSON文件')
    parser.add_argument('--selection_json', help='选型结果JSON（含camera/lens/light）')
    parser.add_argument('--fov', help='视野，如 63x42（mm）')
    parser.add_argument('--project_name', help='项目名称')
    parser.add_argument('--hardware_page', type=int, default=None,
                        help='替换第几个硬件选型页（1起，缺省全部）')

    args = parser.parse_args()

    camera, lens, light_source = {}, {}, {}
    if args.selection_json:
        with open(args.selection_json, 'r', encoding='utf-8') as f:
            sel = json.load(f)
        hw = sel.get('hardware_selection', {})
        camera = hw.get('camera', {})
        lens = hw.get('lens', {})
        light_source = hw.get('light_source', {})
    if args.camera_json:
        with open(args.camera_json, 'r', encoding='utf-8') as f:
            camera = json.load(f)
    if args.lens_json:
        with open(args.lens_json, 'r', encoding='utf-8') as f:
            lens = json.load(f)
    if args.light_json:
        with open(args.light_json, 'r', encoding='utf-8') as f:
            light_source = json.load(f)

    fov = None
    if args.fov:
        parts = args.fov.lower().replace('mm', '').split('x')
        fov = {'width': float(parts[0]), 'height': float(parts[1])}

    generate_ppt(
        template_path=args.template,
        output_path=args.output,
        camera=camera,
        lens=lens,
        light_source=light_source,
        project_name=args.project_name,
        field_of_view=fov,
        hardware_page=args.hardware_page,
    )


if __name__ == '__main__':
    main()
