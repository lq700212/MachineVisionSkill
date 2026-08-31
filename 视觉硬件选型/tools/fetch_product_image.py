#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""官网产品图自动获取：选型后为相机补官方产品图（标准横幅：型号+描述+产品图）

链路（海康API通道，纯HTTP零浏览器，2026-08官网实测）：
  VisionProductType(系列树) → 官网系列名前缀匹配
  → Vision/Cameras?seriesName=xx（型号→productId，productModel 全等匹配防串行）
  → VisionProductIntroduction?id=xx（productModel 二次校验 + previewUrl + 中文描述）
  → 下载产品图（需 referer）→ PIL 合成 662x163 标准横幅 → images/{model}.png
  → 写回库 image_path

诚实原则：任何环节失败 exit 2 诚实报错，绝不编造图片；
型号在官网系列列表中不存在 → 报错（下架/新系列，转人工扩库）；
横幅为"型号+描述+产品图"信息合成，非官网像素级截图（原图来自官网API直链）。

用法：
  python fetch_product_image.py fetch --brand 海康威视 --model MV-CS050-10GM
  python fetch_product_image.py batch                 # 库内缺图相机逐个补图
  python fetch_product_image.py fetch ... --no_db     # 只出图不写库
  python fetch_product_image.py fetch ... --force     # 库中已有图也强制重取

换相机图效果预览（高频场景一条命令，产出预览件、绝不动原方案）：
  python fetch_product_image.py replace --pptx 方案.pptx --model MV-CS060-10GM \
      --brand 海康威视 --auto_variant --export
  # 基础版型号官网不在售时 --auto_variant 自动改用主推变体(-PRO>V5)；
  # --image 可跳过取图直接指定横幅文件；--old-model/--old-image 指定被替换的旧图

硬件页三件套产品图（相机+镜头+光源，按选型结果，各件独立降级）：
  python fetch_product_image.py replace --pptx 方案.pptx --all_parts --export
  # 相机：旧图字节精确匹配；镜头/光源：文字锚定（描述文本框左侧同行的图片）
  # 镜头/光源图三级匹配：库image_path > 型号图{型号}.png > 类型图（assets/part_images/
  #   下中文命名资产如 远心镜头/定焦镜头/环光/条光/同轴光.png，与选型type互包含或
  #   别名命中，内置常用光源/镜头别名，实时可增）；缺图部件 WARN 跳过不阻断
  #   --part lens|light|camera 只换单件

部件类型图资产管理（用户新增图后实时建立关联，不改代码）：
  python fetch_product_image.py map --list                # 看资产与别名
  python fetch_product_image.py map --migrate             # images/类型图迁到assets/part_images/
  python fetch_product_image.py map --add 条光 --alias 条形光 --alias 条形光源
  # 新增类型图：把图（如 条光.png）放进 assets/part_images/ 即自动参与匹配；
  # 选型type写法与文件名对不上时用 --add 登记别名（写aliases.json立即生效）
"""
import argparse
import io
import json
import os
import re
import subprocess
import sys
import time
import urllib.parse
import urllib.request

from PIL import Image, ImageDraw, ImageFont

SKILL_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DEFAULT_DB = os.path.join(SKILL_DIR, 'config', 'hardware_database.json')
IMAGES_DIR = os.path.join(SKILL_DIR, 'images')
if sys.platform == 'win32':
    try:
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
        sys.stderr.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        pass

UA = ('Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 '
      '(KHTML, like Gecko) Chrome/120.0 Safari/537.36')
TIMEOUT_S = 20

# 品牌图片通道适配器（未列出的品牌诚实报"未接入"，绝不假装能取）
BRAND_ADAPTERS = {
    '海康威视': {
        'type': 'hik_api',
        'api_base': 'https://www.hikrobotics.com/cn/Api/Foreground/Vision/',
        'img_base': 'https://www.hikrobotics.com',
        'referer': 'https://www.hikrobotics.com/cn/machinevision/productdetail/',
    },
}

BANNER_W, BANNER_H = 662, 163
BANNER_BG = (232, 237, 243)          # 与 images/ 现有横幅一致的浅蓝灰
FONT_BOLD = r'C:\Windows\Fonts\msyhbd.ttc'
FONT_TEXT = r'C:\Windows\Fonts\msyh.ttc'

# 镜头/光源产品图为本地资产（images/下按型号命名，官网通道未接入的品牌由人工补充）；
# PPT硬件页通过"文字锚定"定位：硬件描述文本框左侧同行的图片即其产品图。
# 左侧光路示意图的部件标签（如"远心镜头"标注框）左侧无照片类图片，天然跳过。
LENS_ANCHORS = ('远心镜头', '镜头')
LIGHT_ANCHORS = ('光源', '环光', '同轴')
IMG_EXTS = ('.png', '.jpg', '.jpeg', '.bmp')

# ── 部件类型图资产（用户中文命名，实时可增）──────────────────────────
# 资产目录：类型代表图（远心镜头.png/环光.png/条光.png...），与相机横幅(images/)分离；
# 兼容扫描 images/（历史类型图未迁移前不断供）。map --migrate 一键搬家。
ASSETS_DIR = os.path.join(SKILL_DIR, 'assets', 'part_images')
ALIASES_JSON = os.path.join(ASSETS_DIR, 'aliases.json')
PART_SCAN_DIRS = (ASSETS_DIR, IMAGES_DIR)

# 内置领域别名表：标准名(=图片文件名去扩展名) → [选型type中的常见写法]。
# aliases.json 可追加/覆盖（map --add 登记），新增类型不改代码。
DEFAULT_ALIASES = {
    '远心镜头': ['双远心', '远心'],
    '定焦镜头': ['定焦'],
    '变焦镜头': ['变焦'],
    '环光': ['环形光', '环形光源', '环型光'],
    '条光': ['条形光', '条形光源', '条型光'],
    '同轴光': ['同轴光源', '同轴'],
    '点光': ['点光源', '点状光'],
    '圆顶光': ['穹顶光', '穹顶光源', '圆顶光源'],
    '背光': ['背光源', '平板背光'],
    '无影光': ['无影光源', '四面无影光'],
    'AOI光': ['AOI光源'],
}


class FetchError(Exception):
    """可报告的获取失败（网络/型号不存在/校验不过）"""


def _http_get(url, referer, binary=False, retries=1):
    """带 UA/referer 的 GET；失败重试1次；返回文本或bytes"""
    last = None
    for _ in range(retries + 1):
        try:
            req = urllib.request.Request(url, headers={
                'User-Agent': UA, 'Referer': referer,
                'Accept': '*/*',
            })
            with urllib.request.urlopen(req, timeout=TIMEOUT_S) as r:
                data = r.read()
                return data if binary else data.decode('utf-8', 'replace')
        except Exception as e:  # noqa: BLE001 - 网络/HTTP错误统一诚实上报
            last = e
            time.sleep(1)
    raise FetchError(f"HTTP请求失败: {url} ({last})")


def _get_json(url, referer):
    text = _http_get(url, referer)
    try:
        return json.loads(text)
    except json.JSONDecodeError:
        raise FetchError(f"接口返回非JSON（可能被反爬拦截）: {url}")


def _require_ok(payload, url):
    """海康API统一包装：status=success 且 data 非空才算拿到"""
    if not isinstance(payload, dict) or payload.get('status') != 'success':
        raise FetchError(f"接口返回error: {url}")
    if payload.get('data') in (None, [], {}):
        raise FetchError(f"接口返回success但data为空: {url}")
    return payload['data']


def resolve_official_series(brand_cfg, series_short):
    """库内短系列名(CS系列) → 官网全名(CS系列工业面阵相机)。
    前缀匹配；0个或多个命中都报错转人工（不猜）。"""
    url = brand_cfg['api_base'] + 'VisionProductType'
    tree = _require_ok(_get_json(url, brand_cfg['referer']), url)
    names = []

    def walk(nodes):
        for n in nodes or []:
            if isinstance(n, dict):
                name = n.get('name') or ''
                if name:
                    names.append(name)
                walk(n.get('children'))

    walk(tree if isinstance(tree, list) else [])
    hits = sorted({n for n in names
                   if n.startswith(series_short) and n != series_short
                   and not n.endswith('全系列')})
    if not hits:
        # 官网可能就叫短名
        hits = sorted({n for n in names if n == series_short})
    if len(hits) != 1:
        raise FetchError(
            f"系列名'{series_short}'在官网系列树中命中{len(hits)}个: {hits or '无'}，转人工确认")
    return hits[0]


def match_model_exact(items, model):
    """productModel 全等匹配（防串行核心：-PRO/-V5/NPOE 变体绝不混淆）"""
    for it in items:
        if isinstance(it, dict) and it.get('productModel') == model:
            return it
    return None


def find_model_variants(items, model):
    """型号基础版不在售时，找官网列表中同前缀在售变体
    （MV-CS060-10GM → MV-CS060-10GM-PRO / MV-CS060-10GM V5，不混入GC/UM色彩接口变体）"""
    return [it['productModel'] for it in items
            if isinstance(it, dict)
            and isinstance(it.get('productModel'), str)
            and it['productModel'].startswith(model)
            and it['productModel'] != model]


def fetch_hik_image(brand_cfg, model, series_short):
    """海康API通道：返回 (product_img PIL, desc 中文描述, product_id)"""
    official_series = resolve_official_series(brand_cfg, series_short)
    print(f"  官网系列名: {official_series}")

    url = (brand_cfg['api_base'] + 'Cameras?seriesName='
           + urllib.parse.quote(official_series))
    items = _require_ok(_get_json(url, brand_cfg['referer']), url)
    hit = match_model_exact(items, model)
    if hit is None:
        variants = find_model_variants(items, model)
        msg = f"型号{model}不在官网系列'{official_series}'列表中（共{len(items)}个型号）"
        if variants:
            msg += f"。基础版可能下架，官网同款在售变体: {' / '.join(variants)}"
        else:
            msg += "——可能下架或属其他系列，转人工"
        raise FetchError(msg)
    pid = hit.get('id')
    print(f"  型号匹配: {model} -> productId={pid}")

    url = brand_cfg['api_base'] + 'VisionProductIntroduction?id=' + str(pid)
    intro = _require_ok(_get_json(url, brand_cfg['referer']), url)
    if intro.get('productModel') != model:
        raise FetchError(
            f"详情返回型号{intro.get('productModel')!r}与目标{model!r}不一致（串行防护拦截）")
    desc = intro.get('productName') or ''
    prev = intro.get('previewUrl') or ''
    if not prev:
        raise FetchError(f"详情无previewUrl产品图: {model}")
    if desc.startswith(model + ',') or desc.startswith(model + '，'):
        desc = desc[len(model) + 1:].strip()
    print(f"  详情校验通过, 描述: {desc[:40]}")

    img_url = brand_cfg['img_base'] + prev
    data = _http_get(img_url, brand_cfg['referer'], binary=True)
    try:
        img = Image.open(io.BytesIO(data))
        img.load()
    except Exception as e:  # noqa: BLE001
        raise FetchError(f"产品图下载后无法解析为图片: {img_url} ({e})")
    if img.width < 100 or img.height < 100:
        raise FetchError(f"产品图尺寸过小({img.size})，疑似占位图: {img_url}")
    print(f"  产品图下载: {img.size} {img.mode}")
    return img, desc, pid


def _load_font(path, size):
    try:
        return ImageFont.truetype(path, size)
    except Exception:  # noqa: BLE001 - 字体缺失降级（中文可能显示为方块）
        return ImageFont.load_default()


def wrap_text(draw, text, font, max_width):
    """中文按字符、英文按单词的宽度断行"""
    lines, cur = [], ''
    for ch in text:
        if draw.textlength(cur + ch, font=font) > max_width and cur:
            lines.append(cur)
            cur = ch
        else:
            cur += ch
    if cur:
        lines.append(cur)
    return lines


def compose_banner(model, desc, product_img, out_path):
    """合成 662x163 标准横幅：左型号名+描述，右产品图（与现有 images 风格一致）"""
    canvas = Image.new('RGB', (BANNER_W, BANNER_H), BANNER_BG)
    d = ImageDraw.Draw(canvas)
    f_model = _load_font(FONT_BOLD, 30)
    f_desc = _load_font(FONT_TEXT, 14)

    d.text((20, 16), model, font=f_model, fill=(25, 25, 25))

    lines = wrap_text(d, desc, f_desc, 350)
    if len(lines) > 3:
        lines = lines[:3] + ['…']  # 超长截断（实际描述均为1~2行）
    y = 76
    for ln in lines:
        d.text((21, y), ln, font=f_desc, fill=(70, 70, 70))
        y += 24

    prod = product_img.convert('RGBA')
    ph = min(148, BANNER_H - 15)
    pw = int(prod.width * ph / prod.height)
    if pw > 260:  # 超宽图限宽再等比
        pw = 260
        ph = int(prod.height * pw / prod.width)
    prod = prod.resize((pw, ph), Image.LANCZOS)
    canvas.paste(prod, (BANNER_W - pw - 14, (BANNER_H - ph) // 2), prod)
    canvas.save(out_path)
    return out_path


def sanitize_filename(model):
    return re.sub(r'[\\/:*?"<>|]', '-', model).strip()


def fetch_camera_image(brand, model, db_path=DEFAULT_DB, update_db=True,
                       out_dir=None, series=None, force=False):
    """主入口：为库内(或指定)相机补官方产品图。
    库中已有本型号产品图且文件在位时直接复用（force=True 强制重取）。
    返回 dict(ok, image_path, message)；失败 ok=False（不抛出，便于批量）。"""
    cfg = BRAND_ADAPTERS.get(brand)
    if cfg is None:
        return {'ok': False, 'message': f"品牌'{brand}'的产品图通道未接入（当前支持: "
                f"{'/'.join(BRAND_ADAPTERS)}），转人工补图"}
    entry = None
    if db_path and os.path.exists(db_path):
        from database_updater import load_db
        db = load_db(db_path)
        entry = next((c for c in db.get('cameras', [])
                      if c.get('model') == model), None)
        if entry is not None and not series:
            series = entry.get('series')
    if entry is not None and not force:
        img_rel = entry.get('image_path') or ''
        img_abs = img_rel if os.path.isabs(img_rel) else os.path.join(SKILL_DIR, img_rel)
        if img_rel and os.path.exists(img_abs):
            print(f"[{model}] 库中已有产品图，直接复用: {img_rel}")
            return {'ok': True, 'image_path': img_abs, 'db_updated': False,
                    'message': '复用库中已有产品图'}
    if not force:
        # 库无条目但本地横幅文件已在（此前取过）→ 同样复用，高频场景不重复请求官网
        local = os.path.join(out_dir or IMAGES_DIR, sanitize_filename(model) + '.png')
        if os.path.exists(local):
            print(f"[{model}] 本地已有产品图，直接复用: {local}")
            return {'ok': True, 'image_path': local, 'db_updated': False,
                    'message': '复用本地已有产品图'}
    if not series:
        # 库无条目/无系列字段：从型号前缀推系列短名（MV-CS050-10GM -> CS系列）
        m = re.match(r'MV-([A-Z]{2})\d', model)
        if not m:
            return {'ok': False, 'message': f"无法从型号'{model}'推断系列，"
                    "请用 --series 指定（如 CS系列）"}
        series = m.group(1) + '系列'

    print(f"[{model}] 品牌={brand} 系列={series}")
    try:
        img, desc, pid = fetch_hik_image(cfg, model, series)
    except FetchError as e:
        print(f"  [FAIL] {e}")
        return {'ok': False, 'message': str(e)}

    out_dir = out_dir or IMAGES_DIR
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, sanitize_filename(model) + '.png')
    compose_banner(model, desc, img, out_path)
    print(f"  横幅已生成: {out_path}")

    rel = os.path.relpath(out_path, SKILL_DIR).replace(os.sep, '/')
    written = False
    if update_db and db_path and os.path.exists(db_path):
        from database_updater import load_db, save_db
        db = load_db(db_path)
        for c in db.get('cameras', []):
            if c.get('model') == model:
                c['image_path'] = rel
                c['source_product_id'] = pid
                written = True
        if written:
            save_db(db, db_path)
            print(f"  库已更新: image_path={rel}")
    return {'ok': True, 'image_path': out_path, 'db_updated': written,
            'message': f"出图成功 desc={desc[:30]}"}


def cmd_fetch(args):
    r = fetch_camera_image(args.brand, args.model, db_path=args.db,
                           update_db=not args.no_db, series=args.series,
                           force=args.force)
    if not r['ok']:
        print(f"结论: ✗ {args.model} 获取失败（{r['message']}）")
        sys.exit(2)
    print(f"结论: ✓ {args.model} 产品图就绪 -> {r['image_path']}")


def cmd_batch(args):
    from database_updater import load_db
    db = load_db(args.db)
    cams = db.get('cameras', [])
    todo = [c for c in cams
            if not c.get('image_path')
            or not os.path.exists(os.path.join(SKILL_DIR, c['image_path']))]
    print(f"缺图相机: {len(todo)}/{len(cams)}")
    fails = []
    for c in todo:
        r = fetch_camera_image(c.get('brand', ''), c['model'], db_path=args.db)
        if not r['ok']:
            fails.append((c['model'], r['message']))
    print(f"\n批处理完成: 成功{len(todo) - len(fails)} 失败{len(fails)}")
    for m, msg in fails:
        print(f"  [未补图] {m}: {msg}")
    if fails:
        sys.exit(2)


def _banner_insert_size(orig_w, orig_h, new_w, new_h):
    """与 generate_ppt.replace_image_in_slide 一致的占位规则：
    保持占位框，按新图宽高比在框内等比缩放"""
    aspect = new_w / new_h
    if orig_w / orig_h > aspect:
        return int(orig_h * aspect), orig_h
    return orig_w, int(orig_w / aspect)


def resolve_old_camera_image(pptx_path, old_model, old_image, db_path):
    """定位PPT中被替换的旧相机图文件。
    优先级：--old-image 直给 > --old-model 查库 > PPT同目录 selection_result.json 选型相机查库。
    定位失败抛 FetchError（诚实转人工，绝不猜）。"""
    if old_image:
        return old_image, '--old-image 指定'
    from database_updater import load_db
    db = load_db(db_path)

    def db_image(m):
        for c in db.get('cameras', []):
            if c.get('model') == m:
                rel = c.get('image_path') or ''
                return rel if os.path.isabs(rel) else os.path.join(SKILL_DIR, rel)
        return None

    if old_model:
        p = db_image(old_model)
        if not p:
            raise FetchError(f"库中无型号'{old_model}'的产品图记录，"
                             "请用 --old-image 直接指定PPT中当前的产品图文件")
        return p, f"--old-model {old_model}(库)"
    sel = os.path.join(os.path.dirname(os.path.abspath(pptx_path)),
                       'selection_result.json')
    if os.path.exists(sel):
        try:
            with open(sel, encoding='utf-8') as fh:
                cam = (json.load(fh).get('hardware_selection') or {}).get('camera') or {}
        except Exception as e:  # noqa: BLE001
            raise FetchError(f"selection_result.json 解析失败: {e}")
        m = cam.get('model')
        p = db_image(m) if m else None
        if p:
            return p, f"同目录selection_result.json选型相机{m}(库)"
    raise FetchError("无法定位PPT中被替换的旧相机图：请用 --old-model 型号 "
                     "或 --old-image 旧图文件路径 指定")


def replace_ppt_camera_image(pptx_path, new_image_path, old_image_path, out_path):
    """把PPT中所有与旧图字节一致的产品图原位替换为新横幅。
    副本输出（先拷贝原文件到 out_path 再改副本），原文件绝不改动。
    返回 (替换明细list, 全页图片清单list)；0命中时清单用于诚实报错。"""
    import shutil
    from pptx import Presentation

    with open(old_image_path, 'rb') as fh:
        old_bytes = fh.read()
    nw, nh = Image.open(new_image_path).size
    shutil.copy2(pptx_path, out_path)
    prs = Presentation(out_path)
    replaced, inventory = [], []
    for si, slide in enumerate(prs.slides, 1):
        for shape in list(slide.shapes):
            if shape.shape_type != 13:  # PICTURE
                continue
            try:
                blob = shape.image.blob
            except Exception:
                continue  # 占位符/链接图无blob
            inventory.append(f"  第{si}页 图片 "
                             f"{shape.width / 914400 * 2.54:.1f}x{shape.height / 914400 * 2.54:.1f}cm "
                             f"pos=({shape.left},{shape.top})")
            if blob != old_bytes:
                continue
            L, T, W, H = shape.left, shape.top, shape.width, shape.height
            new_w, new_h = _banner_insert_size(W, H, nw, nh)
            shape._element.getparent().remove(shape._element)
            slide.shapes.add_picture(new_image_path, L, T, new_w, new_h)
            replaced.append(f"  第{si}页 -> {os.path.basename(new_image_path)} "
                            f"({new_w / 914400 * 2.54:.1f}x{new_h / 914400 * 2.54:.1f}cm)")
    if replaced:
        prs.save(out_path)
    else:
        os.remove(out_path)  # 0命中不留半成品
    return replaced, inventory


def load_aliases(path=None):
    """内置领域别名表 + aliases.json 追加/覆盖（用户实时登记，不改代码）"""
    merged = {k: list(v) for k, v in DEFAULT_ALIASES.items()}
    p = path or ALIASES_JSON
    if os.path.exists(p):
        try:
            with open(p, encoding='utf-8') as fh:
                user = json.load(fh)
            for k, v in (user or {}).items():
                merged.setdefault(k, [])
                for a in v or []:
                    if a not in merged[k]:
                        merged[k].append(a)
        except Exception as e:  # noqa: BLE001 - 配置坏不致命，用内置表
            print(f"  [WARN] aliases.json 解析失败（用内置别名表）: {e}")
    return merged


def scan_part_images(img_dirs=None):
    """扫描资产目录：文件名stem -> 路径（同stem资产目录优先，扩展名不限）"""
    dirs = img_dirs or PART_SCAN_DIRS
    out = {}
    for d in dirs:
        if not os.path.isdir(d):
            continue
        for f in sorted(os.listdir(d)):
            stem, ext = os.path.splitext(f)
            if ext.lower() in IMG_EXTS and stem and stem not in out:
                out[stem] = os.path.join(d, f)
    return out


def match_type_image(part, type_str, img_dirs=None, aliases=None):
    """类型级匹配：资产文件名与选型type互包含，或文件名别名命中type。
    多命中取最具体（stem最长）；返回 (路径, 命中说明) 或 (None, None)。"""
    if not type_str or part not in ('lens', 'light'):
        return None, None
    t = type_str.replace(' ', '').strip()
    if not t:
        return None, None
    stems = scan_part_images(img_dirs)
    alias_map = aliases if aliases is not None else load_aliases()
    hits = []
    for stem, path in stems.items():
        s = stem.replace(' ', '').strip()
        if not s or len(s) > 30:
            continue  # 超长stem基本是误放的说明文件，不参与类型匹配
        if s in t or t in s:
            hits.append((stem, path, '文件名直配'))
        else:
            for a in alias_map.get(stem, []):
                if a in t:
                    hits.append((stem, path, f'别名[{a}]'))
                    break
    if not hits:
        return None, None
    hits.sort(key=lambda x: -len(x[0]))
    stem, path, how = hits[0]
    return path, f"{how}:{stem}"


def migrate_part_images(src_dir=IMAGES_DIR, dst_dir=ASSETS_DIR):
    """把 src_dir 下类型图（非 ASCII 型号命名，如 远心镜头.png）迁到 dst_dir。
    相机横幅（MV-XX 型号样式命名）留在原处。返回迁移明细 list。"""
    os.makedirs(dst_dir, exist_ok=True)
    import shutil
    moved = []
    if not os.path.isdir(src_dir):
        return moved
    for f in sorted(os.listdir(src_dir)):
        stem, ext = os.path.splitext(f)
        if ext.lower() not in IMG_EXTS or not stem:
            continue
        if re.fullmatch(r'[A-Za-z0-9 ()\-\.]+', stem):
            continue  # 型号样式（相机横幅等）不迁移
        dst = os.path.join(dst_dir, f)
        if os.path.exists(dst):
            moved.append(f"  [跳过] {f}（目标已存在同名图）")
            continue
        shutil.move(os.path.join(src_dir, f), dst)
        moved.append(f"  {f} -> {dst_dir}")
    return moved


def find_local_image(model, db_path=DEFAULT_DB, part=None, type_str=None,
                     img_dirs=None):
    """部件产品图查找（三级）：库条目 image_path 显式映射 > 型号图{model}.{ext}
    （扫描目录）> 类型图（资产文件名/别名 与选型type匹配，实时可增）。
    找不到返回 None（诚实降级，绝不就近猜图）。"""
    if db_path and os.path.exists(db_path):
        from database_updater import load_db
        db = load_db(db_path)
        for key in ('cameras', 'lenses', 'light_sources'):
            for c in db.get(key, []):
                if c.get('model') == model:
                    rel = c.get('image_path') or ''
                    p = rel if os.path.isabs(rel) else os.path.join(SKILL_DIR, rel)
                    if rel and os.path.exists(p):
                        return os.path.normpath(p)  # 库内正斜杠归一为系统分隔符
    dirs = img_dirs or PART_SCAN_DIRS
    base = sanitize_filename(model)
    for d in dirs:
        for ext in IMG_EXTS:
            p = os.path.join(d, base + ext)
            if os.path.exists(p):
                return p
    if part and type_str:
        p, _ = match_type_image(part, type_str, img_dirs=dirs)
        if p:
            return p
    return None


def _contain_insert(L, T, W, H, new_w, new_h):
    """原框内等比 contain 居中（任意形状产品图不变形不溢出）"""
    scale = min(W / new_w, H / new_h)
    w, h = int(new_w * scale), int(new_h * scale)
    return L + (W - w) // 2, T + (H - h) // 2, w, h


def anchor_part_shapes(prs, keywords, exclude_shapes=()):
    """文字锚定定位部件产品图：含关键词的文本框左侧同行最近的 PICTURE。
    返回 [(slide_idx, shape, 锚文本)]；锚不到的文本框跳过（如光路示意图部件标签）。
    exclude_shapes: 需排除的 shape id 集合（如已替换过的相机横幅）。"""
    from pptx.util import Emu
    exclude_ids = {id(s._element) for s in exclude_shapes}
    hits = []
    for si, slide in enumerate(prs.slides, 1):
        pics = [s for s in slide.shapes if s.shape_type == 13
                and id(s._element) not in exclude_ids]
        for sh in slide.shapes:
            if not sh.has_text_frame or id(sh._element) in exclude_ids:
                continue
            text = sh.text_frame.text or ''
            if not any(k in text for k in keywords):
                continue
            best, best_gap = None, None
            for p in pics:
                v_overlap = (min(sh.top + sh.height, p.top + p.height)
                             - max(sh.top, p.top))
                if v_overlap <= 0 or p.left >= sh.left:
                    continue  # 不同行或在文本框右侧（右侧为文字排布区约定）
                gap = sh.left - (p.left + p.width)
                if best is None or gap < best_gap:
                    best, best_gap = p, gap
            if best is not None:
                hits.append((si, best, text.strip()[:30]))
    return hits


def replace_part_images(pptx_path, parts, out_path, exclude_shapes=()):
    """把硬件页的镜头/光源产品图按文字锚定原位替换。
    parts = {'lens': 图片路径或None, 'light': 图片路径或None}（None=该部件跳过）。
    副本输出；返回 (明细list, 图片资产清单list, 各部件命中数dict)。"""
    import shutil
    from pptx import Presentation

    shutil.copy2(pptx_path, out_path)
    prs = Presentation(out_path)
    part_names = {'lens': '镜头', 'light': '光源'}
    detail, inventory, hit_count = [], [], {'lens': 0, 'light': 0}
    used = list(exclude_shapes)
    for part in ('lens', 'light'):
        img = parts.get(part)
        if not img:
            continue
        nw, nh = Image.open(img).size
        for si, shape, anchor in anchor_part_shapes(prs, _anchors_for(part), used):
            L, T, W, H = shape.left, shape.top, shape.width, shape.height
            inventory.append(f"  第{si}页 [{part_names[part]}锚:{anchor}] "
                             f"图 {W / 914400 * 2.54:.1f}x{H / 914400 * 2.54:.1f}cm")
            nl, nt, nwd, nht = _contain_insert(L, T, W, H, nw, nh)
            shape._element.getparent().remove(shape._element)
            slide = prs.slides[si - 1]
            new_pic = slide.shapes.add_picture(img, nl, nt, nwd, nht)
            used.append(new_pic)
            hit_count[part] += 1
            detail.append(f"  第{si}页 {part_names[part]} -> {os.path.basename(img)} "
                          f"({nwd / 914400 * 2.54:.1f}x{nht / 914400 * 2.54:.1f}cm)")
    if any(hit_count.values()):
        prs.save(out_path)
    else:
        os.remove(out_path)
    return detail, inventory, hit_count


def _anchors_for(part):
    return LENS_ANCHORS if part == 'lens' else LIGHT_ANCHORS


def resolve_selection_parts(pptx_path, db_path=DEFAULT_DB, img_dirs=None):
    """从PPT同目录 selection_result.json 读选型三件型号 → 本地图路径
    （库image_path > 型号图 > 类型图[资产目录扫描/别名，实时可增]）。
    返回 (parts_dict, notes_list)；缺图部件值为 None 并在 notes 说明。"""
    sel = os.path.join(os.path.dirname(os.path.abspath(pptx_path)),
                       'selection_result.json')
    if not os.path.exists(sel):
        return {}, ["未找到 selection_result.json（请先跑主流程生成，或用 --model/--image 显式指定）"]
    with open(sel, encoding='utf-8') as fh:
        hs = json.load(fh).get('hardware_selection') or {}
    parts, notes = {}, []
    for part, key in (('camera', 'camera'), ('lens', 'lens'), ('light', 'light_source')):
        ent = hs.get(key) or {}
        model = ent.get('model')
        if not model:
            notes.append(f"  [{key}] 选型结果无型号，跳过")
            parts[part] = None
            continue
        img = find_local_image(model, db_path, part=part,
                               type_str=ent.get('type') or '', img_dirs=img_dirs)
        parts[part] = img
        if img:
            via = ('类型图' if os.path.splitext(os.path.basename(img))[0]
                   not in (model, sanitize_filename(model)) else '型号图')
            notes.append(f"  [{key}] {model} -> {os.path.basename(img)}（{via}）")
        else:
            hint = ''
            if part in ('lens', 'light'):
                hit, _ = match_type_image(part, ent.get('type') or '', img_dirs=img_dirs)
                if hit is None:
                    have = [s for s in scan_part_images(img_dirs)]
                    if have:
                        hint = (f"；资产目录现有 {('/'.join(have[:8]))}，"
                                f"若含对应类型图请 map --add {have[0]} --alias <type关键词> 登记")
            notes.append(f"  [{key}] {model} 无匹配产品图{hint}，该部件跳过")
    return parts, notes


def _export_pages(out_path):
    """导出页面图到同目录 preview_check（失败不阻断）"""
    outdir = os.path.join(os.path.dirname(os.path.abspath(out_path)), 'preview_check')
    print(f"[导出] 页面图 -> {outdir}")
    try:
        subprocess.run(
            [sys.executable,
             os.path.join(os.path.dirname(os.path.abspath(__file__)),
                          'export_ppt_images.py'), out_path, '--outdir', outdir],
            check=True, timeout=300)
    except Exception as e:  # noqa: BLE001
        print(f"  [WARN] 页面图导出失败（不影响已产出的PPT）: {e}")


def _cmd_replace_parts(args):
    """按选型结果替换硬件页产品图：相机(blob精确匹配) + 镜头/光源(文字锚定)。
    各部件独立降级：缺图/未锚定均 WARN 不阻断，绝不瞎猜就近换图。"""
    import shutil
    if not os.path.exists(args.pptx):
        print(f"结论: ✗ PPT不存在: {args.pptx}")
        sys.exit(2)
    parts, notes = resolve_selection_parts(args.pptx, args.db)
    print("[选型部件图片解析]")
    print('\n'.join(notes) or '  （无）')
    if args.part:  # 单件模式：只启用指定部件
        parts = {k: (v if k == args.part else None) for k, v in parts.items()}

    stem = os.path.splitext(os.path.basename(args.pptx))[0]
    out = args.out or os.path.join(os.path.dirname(os.path.abspath(args.pptx)),
                                   f"{stem}_硬件图预览.pptx")
    tmp_a, tmp_b = out + '.a.tmp', out + '.b.tmp'  # 两级交替中转，杜绝自拷自
    for f in (tmp_a, tmp_b):
        if os.path.exists(f):
            os.remove(f)
    cur_src = args.pptx
    summary = []

    # 1) 相机：旧图字节精确匹配
    if parts.get('camera'):
        try:
            old_img, how = resolve_old_camera_image(args.pptx, None, None, args.db)
        except FetchError as e:
            old_img, how = None, str(e)
        if old_img and os.path.exists(old_img):
            rep, _ = replace_ppt_camera_image(cur_src, parts['camera'], old_img,
                                              tmp_a)
            if rep:
                summary.append(f"相机{len(rep)}处")
                cur_src = tmp_a
            else:
                summary.append(f"相机0处（旧图未命中: {how}）")
        else:
            summary.append(f"相机跳过（{how}）")
    else:
        summary.append("相机跳过（无本地图）")

    # 2) 镜头/光源：文字锚定（输入输出恒为不同文件）
    part_imgs = {'lens': parts.get('lens'), 'light': parts.get('light')}
    if any(part_imgs.values()):
        detail, _, hits = replace_part_images(cur_src, part_imgs, tmp_b)
        for part, cn in (('lens', '镜头'), ('light', '光源')):
            if hits.get(part):
                summary.append(f"{cn}{hits[part]}处")
                cur_src = tmp_b
            else:
                summary.append(f"{cn}0处" + ("（无本地图）" if not part_imgs.get(part)
                                             else "（页面未锚定到图片，转人工确认布局）"))
        if detail:
            print("[替换明细]")
            print('\n'.join(detail))
    else:
        summary.append("镜头/光源跳过（无本地图）")

    if cur_src == args.pptx:
        for f in (tmp_a, tmp_b):
            if os.path.exists(f):
                os.remove(f)
        print("结论: ✗ 所有部件均未替换（0处），未产出文件。"
              "请确认 images/ 资产命名与硬件页布局")
        sys.exit(2)
    if cur_src != out:
        shutil.move(cur_src, out)
    for f in (tmp_a, tmp_b):
        if os.path.exists(f):
            os.remove(f)
    print("[汇总] " + "；".join(summary))
    print(f"[输出] {out}（原文件未改动）")
    if args.export:
        _export_pages(out)
    print(f"结论: ✓ 硬件图预览就绪 -> {out}")


def cmd_replace(args):
    if getattr(args, 'all_parts', False) or getattr(args, 'part', None):
        _cmd_replace_parts(args)
        return
    if not args.image and not args.model:
        print("结论: ✗ 需要 --model 型号（官网取图）或 --image 文件（直接指定横幅）")
        sys.exit(2)

    # 1) 新图就绪：--image 直给 > 库中已有复用 > 官网取图（可自动变体）
    if args.image:
        if not os.path.exists(args.image):
            print(f"结论: ✗ --image 文件不存在: {args.image}")
            sys.exit(2)
        new_img = args.image
        print(f"[新图] 直接使用指定文件: {new_img}")
    else:
        r = fetch_camera_image(args.brand, args.model, db_path=args.db)
        if not r['ok'] and args.auto_variant:
            m = re.search(r'在售变体: (.+)$', r['message'])
            if m:
                cands = [x.strip() for x in m.group(1).split('/') if x.strip()]
                cands.sort(key=lambda x: (not x.endswith('-PRO'), ' V5' not in x))
                print(f"  [自动变体] 基础版不在售，改用主推变体 {cands[0]}"
                      f"（其余: {' / '.join(cands[1:]) or '无'}）")
                args.model = cands[0]
                r = fetch_camera_image(args.brand, args.model, db_path=args.db)
        if not r['ok']:
            print(f"结论: ✗ 新图获取失败（{r['message']}）")
            sys.exit(2)
        new_img = r['image_path']

    # 2) 旧图定位
    try:
        old_img, how = resolve_old_camera_image(args.pptx, args.old_model,
                                                args.old_image, args.db)
    except FetchError as e:
        print(f"结论: ✗ {e}")
        sys.exit(2)
    if not os.path.exists(old_img):
        print(f"结论: ✗ 旧相机图文件不存在: {old_img}（请用 --old-image 指定）")
        sys.exit(2)
    print(f"[旧图] {how}: {old_img}")

    # 3) 替换（副本输出，原文件不动）
    stem = os.path.splitext(os.path.basename(args.pptx))[0]
    out = args.out or os.path.join(os.path.dirname(os.path.abspath(args.pptx)),
                                   f"{stem}_相机图预览_{sanitize_filename(args.model or '指定图')}.pptx")
    replaced, inventory = replace_ppt_camera_image(args.pptx, new_img, old_img, out)
    if not replaced:
        print("结论: ✗ PPT中未找到与旧图一致的产品图（0处替换），未产出文件。"
              "各页图片清单:\n" + ('\n'.join(inventory) or '  （无图片）'))
        print("  → 若原PPT曾被PowerPoint重存（图片字节变化），请用 --old-image 指定其实际图文件")
        sys.exit(2)
    print(f"[替换] 共{len(replaced)}处:")
    print('\n'.join(replaced))
    print(f"[输出] {out}（原文件未改动）")

    # 4) 页面图导出（失败不阻断，PPT已产出）
    if args.export:
        _export_pages(out)
    print(f"结论: ✓ {args.model or '指定图'} 已入预览件（共{len(replaced)}处）-> {out}")


def cmd_map(args):
    """部件类型图资产管理：--list 查看 / --migrate 搬家 / --add 登记别名（实时建立关联）"""
    os.makedirs(ASSETS_DIR, exist_ok=True)
    if args.migrate:
        moved = migrate_part_images()
        print(f"[迁移] images/ -> {ASSETS_DIR}")
        print('\n'.join(moved) or '  无可迁移的类型图（型号横幅不受影响）')
        print(f"结论: ✓ 迁移完成（相机横幅仍留在 images/）")
        return
    if args.list:
        stems = scan_part_images()
        alias_map = load_aliases()
        print(f"[资产目录] {ASSETS_DIR}")
        print(f"[兼容扫描] {IMAGES_DIR}")
        if not stems:
            print("  （无图片。类型图放此处即自动参与匹配，如 远心镜头.png/条光.png）")
        for stem, p in stems.items():
            al = alias_map.get(stem, [])
            print(f"  {os.path.basename(p)}  别名: {'/'.join(al) or '无'}")
        print(f"结论: ✓ 共{len(stems)}张部件图")
        return
    if args.add and args.alias:
        stems = scan_part_images()
        if args.add not in stems:
            print(f"结论: ✗ 资产目录不存在图'{args.add}'（{ASSETS_DIR}\\{args.add}.png/.jpg），"
                  "拒绝登记空映射——先放图再登记")
            sys.exit(2)
        user = {}
        if os.path.exists(ALIASES_JSON):
            try:
                with open(ALIASES_JSON, encoding='utf-8') as fh:
                    user = json.load(fh)
            except Exception:
                user = {}
        lst = user.setdefault(args.add, [])
        added = [a for a in args.alias if a != args.add and a not in lst]
        if not added:
            print("结论: ✓ 别名均已登记，无变更")
            return
        lst.extend(added)
        with open(ALIASES_JSON, 'w', encoding='utf-8') as fh:
            json.dump(user, fh, ensure_ascii=False, indent=2)
        print(f"结论: ✓ 已登记 {args.add} <- {'/'.join(added)}（写入aliases.json，立即生效）")
        return
    print("结论: ✗ map 需要 --list / --migrate / --add 标准名 --alias 别名 之一")
    sys.exit(2)


def main():
    ap = argparse.ArgumentParser(description='官网产品图自动获取')
    sub = ap.add_subparsers(dest='cmd', required=True)

    f = sub.add_parser('fetch', help='为指定相机补官方产品图')
    f.add_argument('--brand', required=True)
    f.add_argument('--model', required=True)
    f.add_argument('--series', help='官网系列短名（缺省从型号/库推断）')
    f.add_argument('--db', default=DEFAULT_DB)
    f.add_argument('--no_db', action='store_true', help='只出图不写库')
    f.add_argument('--force', action='store_true',
                   help='库中已有产品图时也强制重新取图')
    f.set_defaults(func=cmd_fetch)

    b = sub.add_parser('batch', help='库内缺图相机逐个补图')
    b.add_argument('--db', default=DEFAULT_DB)
    b.set_defaults(func=cmd_batch)

    m = sub.add_parser('map', help='部件类型图资产管理（扫描/迁移/登记别名，实时建立关联）')
    m.add_argument('--list', action='store_true', help='列出资产图与已登记别名')
    m.add_argument('--migrate', action='store_true',
                   help='把images/下类型图（中文命名等）迁到assets/part_images/，型号横幅不动')
    m.add_argument('--add', help='标准名（=资产图文件名去扩展名，如 条光）')
    m.add_argument('--alias', action='append',
                   help='登记到该标准名的type关键词别名（可多次，如 --alias 条形光 --alias 条形光源）')
    m.set_defaults(func=cmd_map)

    r = sub.add_parser('replace',
                       help='把PPT中的相机产品图换成指定型号官网图（预览件，不动原文件）')
    r.add_argument('--pptx', required=True, help='方案PPT路径')
    r.add_argument('--model', help='新型号（配合 --brand 官网取图）')
    r.add_argument('--image', help='直接指定新横幅文件（跳过官网取图）')
    r.add_argument('--brand', default='海康威视')
    r.add_argument('--old-model', help='被替换的旧相机型号（查库定位旧图；缺省读PPT同目录selection_result.json）')
    r.add_argument('--old-image', help='直接指定被替换的旧图文件路径')
    r.add_argument('--out', help='输出预览件路径（缺省 原名_相机图预览_型号.pptx；--all_parts 时为 原名_硬件图预览.pptx）')
    r.add_argument('--all_parts', action='store_true',
                   help='按选型结果替换硬件页全部产品图（相机+镜头+光源，各件独立降级）')
    r.add_argument('--part', choices=['camera', 'lens', 'light'],
                   help='只替换指定部件（配合选型结果自动取图）')
    r.add_argument('--auto_variant', action='store_true',
                   help='基础版型号官网不在售时自动改用主推在售变体（-PRO > V5 > 其他）')
    r.add_argument('--export', action='store_true', help='替换后自动导出页面图')
    r.add_argument('--db', default=DEFAULT_DB)
    r.set_defaults(func=cmd_replace)

    args = ap.parse_args()
    args.func(args)


if __name__ == '__main__':
    main()
