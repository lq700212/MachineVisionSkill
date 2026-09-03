#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成PPT规则自查工具 - 把视觉验收经验固化为可重复检查
====================================================

背景：视觉验收曾发现4类问题（光源类型矛盾/WD数值断行/标注文字溢出/旧型号残留），
均由评审agent多轮发现。本工具把这些问题的**判定规则**固化，任何模型跑一次脚本
即可得到同等结论的机检报告，不再依赖强模型临场评审。

检查项：
  C1 旧型号残留     硬件选型页不得残留模板旧硬件（FL-COX2/FL-R-/acA2500/Basler等）
  C2 光源类型一致   页面类型词与选中光源类型同义，两硬件页之间无矛盾
  C3 两页参数一致   相机/镜头型号、倍率、像素精度、视野在两个硬件页数值一致
  C4 文本溢出预检   估算文本渲染宽度vs文本框宽度：超宽且自动换行→断行FAIL；
                    超宽且不换行→单行溢出WARN
  C5 数值合理性复算 实际视野≥需求视野、物方分辨率≤像素精度上限、倍率在窗口内
  C6 阴影效果检查   任一页XML含effectLst/outerShdw即FAIL（曾被动过阴影）
  C7 相机图替换     硬件页右上半区应存在产品图

用法：
  python check_ppt_quality.py --pptx 方案.pptx --selection output/selection_result.json
  python check_ppt_quality.py --pptx 方案.pptx --selection ... --out 验收报告.txt

退出码：0=无FAIL（可交付）；1=存在FAIL（按报告处理）；2=输入错误
"""

import argparse
import json
import math
import os
import re
import sys
import zipfile

from pptx import Presentation
from pptx.util import Emu

PT_PER_MM = 2.83465
EMU_PER_MM = 36000.0

# 模板中的旧硬件触发词（出现即说明替换遗漏）
STALE_TOKENS = [
    'FL-COX2', 'FL-R-', 'acA2500', 'Basler', 'DTCM110-80H',
    'FL-CL', 'FL-RL',
]

# 光源类型同义组（同组=语义一致）
LIGHT_TYPE_SYNONYMS = [{'环形光', '环光', '圆环光'}, {'同轴光'}, {'背光'}, {'点光', '点光源'}]

# 检测的硬件页类型词（独立标注框或行首前缀）
LIGHT_TYPE_PATTERN = re.compile(r'(同轴光|环光|环形光|背光|点光源|点光)')


def is_hardware_slide(slide):
    """硬件选型页判定（与generate_ppt一致：短标题为硬件选型类）"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        t = shape.text_frame.text.strip()
        if t and len(t) <= 12 and ('硬件选型' in t or '相机选型' in t or '相机参数' in t):
            return True
        if re.match(r'^(同轴光|环光|环形光|背光|光源)[:：]', t):
            return True
    return False


def slide_texts(slide):
    """收集页面全部段落文本 [(text, shape, para)]"""
    out = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                out.append((t, shape, para))
    return out


def norm_light_type(s):
    """类型词规范化：'LED环形光光源'→'环形光'"""
    s = re.sub(r'^LED|光源$', '', s.strip()).strip()
    return s


def types_consistent(a, b):
    """两个类型词是否同义"""
    if a == b:
        return True
    for group in LIGHT_TYPE_SYNONYMS:
        if a in group and b in group:
            return True
    return False


def est_text_width_pt(text, font_size_pt):
    """估算文本渲染宽度(pt)：中文/全角≈1.0em，ASCII≈0.55em"""
    w = 0.0
    for ch in text:
        if ord(ch) > 0x2E7F:  # CJK/全角
            w += 1.0
        elif ch in '0123456789':
            w += 0.58
        elif ch.isupper():
            w += 0.68
        else:
            w += 0.52
    return w * font_size_pt


def para_font_size(para):
    """取段落第一个明确字号(pt)；继承主题时返回None（无法可靠估算宽度）"""
    for run in para.runs:
        if run.font.size is not None:
            return run.font.size.pt
    return None


def check_stale_models(hw_slides):
    """C1: 旧型号残留"""
    issues = []
    for idx, slide in hw_slides:
        for t, _, _ in slide_texts(slide):
            for token in STALE_TOKENS:
                if token in t:
                    issues.append(f"第{idx}页残留旧内容 {token!r}: {t[:40]!r}")
    return issues


def collect_light_types(hw_slides):
    """收集各硬件页出现的光源类型词"""
    found = {}  # idx -> set(类型词)
    for idx, slide in hw_slides:
        s = set()
        for t, _, _ in slide_texts(slide):
            # 行首带冒号的类型行 或 独立类型标注框
            m = re.match(r'^(同轴光|环光|环形光|背光|LED环形光光源|LED同轴光光源)[:：]', t)
            if m:
                s.add(norm_light_type(m.group(1)))
                continue
            m2 = re.fullmatch(r'(同轴光|环光|环形光|背光)', t)
            if m2:
                s.add(norm_light_type(m2.group(1)))
        found[idx] = s
    return found


def check_light_consistency(hw_slides, sel_light_type):
    """C2: 光源类型一致性（页面vs选型、两页之间）"""
    issues = []
    found = collect_light_types(hw_slides)
    sel_norm = norm_light_type(sel_light_type or '')
    for idx, types in found.items():
        if not types:
            issues.append(f"[WARN] 第{idx}页未识别到光源类型标注（若该页无光源图示可忽略）")
            continue
        for t in types:
            if sel_norm and not types_consistent(t, sel_norm):
                issues.append(
                    f"第{idx}页光源类型 {t!r} 与选中光源 {sel_light_type!r} 矛盾")
    # 两页之间
    idxs = sorted(found.keys())
    for i in range(len(idxs) - 1):
        a, b = found[idxs[i]], found[idxs[i + 1]]
        for ta in a:
            for tb in b:
                if not types_consistent(ta, tb):
                    issues.append(
                        f"第{idxs[i]}页类型 {ta!r} 与第{idxs[i+1]}页类型 {tb!r} 相互矛盾")
    return issues


def extract_features(texts):
    """从页面文本提取关键数值特征"""
    f = {'camera_model': None, 'lens_model': None, 'mag': None,
         'pixel_precision': None, 'fov': None}
    for t, _, _ in texts:
        if f['camera_model'] is None:
            m = re.search(r'(MV-CS\d+-\d+[A-Z]{2}|MV-C[A-Z]\d+-\d+\S*)', t)
            if m:
                f['camera_model'] = m.group(1)
        if f['lens_model'] is None:
            m = re.search(r'(DTCM\d+-\d+\S*|WWK\d+-\d+\S*)', t)
            if m:
                f['lens_model'] = m.group(1)
        if f['mag'] is None:
            m = re.search(r'(\d\.\d+)\s*x\b', t)
            if m:
                f['mag'] = float(m.group(1))
        if f['pixel_precision'] is None:
            m = re.search(r'(\d\.\d+)\s*mm/pixel', t)
            if m:
                f['pixel_precision'] = float(m.group(1))
        if f['fov'] is None:
            m = re.search(r'(\d{2,4})mm\s*[×x]\s*(\d{2,4})mm', t)
            if m:
                f['fov'] = (int(m.group(1)), int(m.group(2)))
    return f


def check_two_page_consistency(hw_slides):
    """C3: 两个硬件页关键参数一致"""
    issues = []
    if len(hw_slides) < 2:
        return issues
    feats = [extract_features(slide_texts(s)) for _, s in hw_slides]
    a, b = feats[0], feats[1]
    for key, name in (('camera_model', '相机型号'), ('lens_model', '镜头型号'),
                      ('mag', '镜头倍率'), ('pixel_precision', '像素精度'),
                      ('fov', '视野')):
        va, vb = a.get(key), b.get(key)
        if va is not None and vb is not None and va != vb:
            issues.append(f"两硬件页{name}不一致: {va} vs {vb}")
    return issues


def check_overflow(hw_slides):
    """C4: 文本溢出预检（把'51.1m/m断行'、'LED环形光光源溢出'规则化）"""
    issues = []
    for idx, slide in hw_slides:
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.width is None:
                continue
            box_pt = Emu(shape.width).pt
            try:
                wrap = shape.text_frame.word_wrap
            except Exception:
                wrap = None
            wrap_off = (wrap is False)
            for para in shape.text_frame.paragraphs:
                t = para.text.strip()
                if not t or len(t) > 60:
                    continue
                fs = para_font_size(para)
                if fs is None:
                    # 字号继承主题：渲染字号未知，仅严重超长才提示人工确认
                    est = est_text_width_pt(t, 12.0)
                    if est > box_pt * 1.8:
                        issues.append(f"[WARN] 第{idx}页 {t[:24]!r} 文本明显偏长"
                                      f"（字号继承主题无法精估），建议人工确认显示完整")
                    continue
                est = est_text_width_pt(t, fs)
                if est > box_pt * 1.25:
                    if wrap_off:
                        issues.append(f"[WARN] 第{idx}页 {t[:24]!r} 估算宽{est:.0f}pt"
                                      f">框宽{box_pt:.0f}pt（不换行，可能压边，建议人工确认）")
                    else:
                        issues.append(
                            f"第{idx}页 {t[:24]!r} 估算宽{est:.0f}pt>框宽{box_pt:.0f}pt"
                            f"且自动换行开启→断行/溢出风险（参照历史缺陷：数值断行、标注溢出）")
                elif est > box_pt * 1.05:
                    issues.append(f"[WARN] 第{idx}页 {t[:24]!r} 文本接近框宽"
                                  f"（{est:.0f}/{box_pt:.0f}pt），留意不同字体下的溢出")
    return issues


def check_values(selection):
    """C5: 数值合理性复算（与需求口径对照）"""
    issues = []
    hs = selection.get('hardware_selection', {})
    perf = selection.get('performance', {})
    pi = selection.get('project_info', {})
    req_fov = hs.get('required_fov') or {}
    act_fov = hs.get('actual_fov') or {}
    lens = hs.get('lens', {})
    if req_fov and act_fov:
        if act_fov.get('width', 0) < req_fov.get('width', 1e9) - 0.05 or \
           act_fov.get('height', 0) < req_fov.get('height', 1e9) - 0.05:
            issues.append(
                f"实际视野{act_fov.get('width'):.1f}x{act_fov.get('height'):.1f}mm "
                f"未覆盖需求{req_fov.get('width'):.1f}x{req_fov.get('height'):.1f}mm")
    obj_res = perf.get('object_resolution_um')
    ppm = pi.get('pixel_precision_max_mm')
    if obj_res is not None and ppm is not None and obj_res > ppm * 1000 + 0.01:
        issues.append(
            f"物方分辨率{obj_res:.2f}μm 超过像素精度上限{ppm*1000:.2f}μm")
    mag = lens.get('magnification')
    win = perf.get('mag_window') or {}
    if mag is not None and win.get('min') is not None:
        if not (win['min'] * 0.98 <= mag <= win['max'] * 1.02):
            issues.append(f"镜头倍率{mag}不在可行窗口[{win['min']}, {win['max']}]")
    if not hs.get('fov_satisfied', True):
        issues.append("选型结果标记 fov_satisfied=False")
    return issues


def check_shadow(pptx_path):
    """C6: 阴影效果检查（曾发生阴影污染事故，持续设防）"""
    issues = []
    try:
        with zipfile.ZipFile(pptx_path) as z:
            for name in z.namelist():
                if not re.match(r'ppt/(slides|slideLayouts|slideMasters)/[^/]+\.xml$', name):
                    continue
                xml = z.read(name).decode('utf-8', errors='ignore')
                n_shdw = xml.count('outerShdw')
                if n_shdw:
                    issues.append(f"{os.path.basename(name)} 含 {n_shdw} 处outerShdw阴影效果")
    except Exception as e:
        issues.append(f"[WARN] 阴影检查未完成: {e}")
    return issues


def check_camera_image(hw_slides, slide_width_emu, slide_height_emu):
    """C7: 硬件页右上半区应有产品图"""
    issues = []
    for idx, slide in hw_slides:
        ok = False
        for shape in slide.shapes:
            if shape.shape_type == 13 and shape.left is not None:  # PICTURE
                if shape.left > slide_width_emu * 0.45 and shape.top < slide_height_emu * 0.55:
                    ok = True
                    break
        if not ok:
            issues.append(f"[WARN] 第{idx}页右上半区未发现产品图片（相机图可能未替换）")
    return issues


def main():
    parser = argparse.ArgumentParser(description='生成PPT规则自查（视觉验收经验固化）')
    parser.add_argument('--pptx', required=True, help='生成的PPTX路径')
    parser.add_argument('--selection', help='selection_result.json路径（缺省取pptx同目录）')
    parser.add_argument('--out', help='验收报告输出路径（缺省只打印）')
    args = parser.parse_args()

    pptx_path = os.path.abspath(args.pptx)
    if not os.path.exists(pptx_path):
        print(f"文件不存在: {pptx_path}")
        return 2
    sel_path = args.selection or os.path.join(os.path.dirname(pptx_path),
                                              'selection_result.json')
    selection = {}
    if os.path.exists(sel_path):
        with open(sel_path, encoding='utf-8') as f:
            selection = json.load(f)
    else:
        print(f"[WARN] 未找到选型结果 {sel_path}，C5数值复算跳过")

    prs = Presentation(pptx_path)
    sw, sh = prs.slide_width, prs.slide_height
    hw_slides = [(i + 1, s) for i, s in enumerate(prs.slides)
                 if is_hardware_slide(s)]
    if not hw_slides:
        print("未识别到硬件选型页（标题含'硬件选型'等），无法自查")
        return 2

    sel_light_type = (selection.get('hardware_selection', {})
                      .get('light_source', {}).get('type', ''))

    results = []  # (级别, 消息)
    c1 = check_stale_models(hw_slides)
    results += [('FAIL', m) for m in c1]
    c2 = check_light_consistency(hw_slides, sel_light_type)
    results += [('FAIL' if not m.startswith('[WARN]') else 'WARN', m) for m in c2]
    c3 = check_two_page_consistency(hw_slides)
    results += [('FAIL', m) for m in c3]
    c4 = check_overflow(hw_slides)
    results += [('FAIL' if not m.startswith('[WARN]') else 'WARN', m) for m in c4]
    c5 = check_values(selection)
    results += [('FAIL', m) for m in c5]
    c6 = check_shadow(pptx_path)
    results += [('FAIL' if not m.startswith('[WARN]') else 'WARN', m) for m in c6]
    c7 = check_camera_image(hw_slides, sw, sh)
    results += [('WARN', m) for m in c7]

    n_fail = sum(1 for lv, _ in results if lv == 'FAIL')
    n_warn = sum(1 for lv, _ in results if lv == 'WARN')
    lines = []
    lines.append("=" * 66)
    lines.append("生成PPT规则自查报告（check_ppt_quality）")
    lines.append("=" * 66)
    lines.append(f"PPT: {os.path.basename(pptx_path)}  硬件页: {[i for i, _ in hw_slides]}")
    lines.append(f"结论: {'✓ 通过（可交付）' if n_fail == 0 else '✗ 存在FAIL（处理后交付）'}"
                 f"  FAIL={n_fail}  WARN={n_warn}")
    lines.append("-" * 66)
    if results:
        for lv, m in results:
            lines.append(f"[{lv}] {m}")
    else:
        lines.append("全部检查项通过：无旧型号残留、光源类型两页一致、两页参数一致、"
                     "无文本溢出风险、数值复算正确、无阴影、相机图已替换")
    lines.append("-" * 66)
    lines.append("故障处理对照：")
    lines.append("  旧型号残留  → 该文本未被替换规则命中，检查generate_ppt触发词")
    lines.append("  光源类型矛盾 → 检查light_selector的type与模板标注框（3.3b同义词逻辑）")
    lines.append("  断行/溢出   → 值变长导致：WD标注需word_wrap=False；短标注框只写短词")
    lines.append("  数值不合理  → 重跑选型（vision_proposal_generator --config ... --auto）")
    report = "\n".join(lines)
    # 根因注释：报告必须先落盘再打印——此前先 print(report)，
    # gbk 控制台遇到 ✓/✗ 直接 UnicodeEncodeError 崩溃，--out 文件没写、
    # 主流程验收误报"存在FAIL项"。且打印本身做容错，避免刷屏 traceback
    if args.out:
        with open(args.out, 'w', encoding='utf-8') as f:
            f.write(report + "\n")
        print(f"\n报告已保存: {args.out}")
    try:
        print(report)
    except UnicodeEncodeError:
        safe = report.replace('✓', 'PASS').replace('✗', 'FAIL')
        print(safe.encode('gbk', errors='replace').decode('gbk'))
    return 0 if n_fail == 0 else 1


if __name__ == '__main__':
    sys.exit(main())
