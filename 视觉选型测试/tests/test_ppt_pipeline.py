# -*- coding: utf-8 -*-
"""端到端重用例：config→选型→核验→PPT→验收FAIL=0
（config 与模板均在临时目录现场合成，不依赖用户工作区、不污染任何真实数据；
2026-09-03 前本文件依赖工作区 project_config.json+模板缺失即整类跳过，
改写为自包含后任何环境都实跑）"""
import glob
import json
import os
import subprocess
import sys
import unittest

import vt_common

GENERATOR = os.path.join(vt_common.TOOLS_DIR, 'vision_proposal_generator.py')

# 标准口径：此前 --text 实跑 FAIL=0 的组合（精度0.03/节拍3s/检测区域38.5x22）
CONFIG = {
    'project_name': '自包含端到端',
    'precision_requirement': 0.03,
    'cycle_time': 3,
    'detection_area': '38.5x22',
}


def write_config(path):
    with open(path, 'w', encoding='utf-8') as f:
        json.dump(CONFIG, f, ensure_ascii=False)


def write_template(path):
    """合成最小硬件模板：一页硬件选型页，含表格4行+描述3行+WD标注2个。
    文本刻意避开 STALE_TOKENS（旧型号残留检查的一切触发词），
    文本框给足宽度并显式字号（溢出预检不误报），无图片（C7仅WARN可接受）。"""
    from pptx import Presentation
    from pptx.util import Mm, Pt
    prs = Presentation()
    prs.slide_width, prs.slide_height = Mm(254), Mm(142)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白版式

    def add_text(left, top, width, height, text, size=14):
        box = slide.shapes.add_textbox(Mm(left), Mm(top), Mm(width), Mm(height))
        box.text_frame.word_wrap = True
        p = box.text_frame.paragraphs[0]
        p.text = text
        p.runs[0].font.size = Pt(size)
        return box

    add_text(10, 5, 200, 12, '硬件选型', size=20)
    # 表格：行首列标题命中替换规则，值列给占位（会被实际值替换）
    tbl_shape = slide.shapes.add_table(4, 2, Mm(10), Mm(20), Mm(200), Mm(60))
    rows = [('相机工作距离', '110mm'), ('相机视野', '50mm×40mm'),
            ('相机精度', '0.01mm/pixel'), ('光源工作距离', '80mm')]
    for i, (h, v) in enumerate(rows):
        for j, t in enumerate((h, v)):
            cell = tbl_shape.table.cell(i, j)
            cell.text = t
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
    add_text(10, 85, 200, 12, '使用500万像素相机，分辨率为2448 × 2048')
    add_text(10, 99, 200, 12, '镜头： WWK05-110-111V3')
    add_text(10, 113, 200, 12, '环光： RL-100-90-W')
    add_text(220, 20, 25, 10, '110mm')
    add_text(220, 35, 25, 10, '80mm')
    prs.save(path)


def run_generator(cfg, template, out, cwd=None):
    cmd = [sys.executable, GENERATOR, '--config', cfg,
           '--template', template, '--output', out, '--auto']
    return subprocess.run(cmd, capture_output=True, text=True,
                          encoding='utf-8', errors='replace',
                          timeout=300, cwd=cwd or os.getcwd())


def assert_deliverable(testcase, proc, out):
    combined = (proc.stdout or '') + (proc.stderr or '')
    testcase.assertEqual(proc.returncode, 0,
                         f'主流程退出码{proc.returncode}\n{combined[-1500:]}')
    pptx = glob.glob(os.path.join(out, '*.pptx'))
    testcase.assertTrue(pptx, f'未产出PPT\n{combined[-800:]}')
    report = os.path.join(out, 'acceptance_report.txt')
    testcase.assertTrue(os.path.exists(report), '验收报告未生成')
    content = open(report, encoding='utf-8', errors='replace').read()
    testcase.assertRegex(content, r'FAIL=0',
                         msg=f'PPT规则自查存在FAIL，不可交付\n{content[-800:]}')


class TestPptPipeline(unittest.TestCase):

    def test_config_to_ppt_fail0(self):
        """回归: config一条命令到PPT，验收FAIL=0（可交付标准）"""
        tmp = vt_common.temp_dir('vt_ppt_')
        try:
            cfg = os.path.join(tmp, 'project_config.json')
            tpl = os.path.join(tmp, 'template.pptx')
            out = os.path.join(tmp, 'out')
            write_config(cfg)
            write_template(tpl)
            p = run_generator(cfg, tpl, out)
            assert_deliverable(self, p, out)
        finally:
            vt_common.rmtree(tmp)

    def test_relative_output_acceptance_still_fails0(self):
        """回归: 相对 --output 时验收不得静默失效（2026-08-31事故：
        验收子进程 cwd=tools_dir，相对路径被 abspath 错位到 tools\\output，
        验收查错文件仍打印完成——交付失检。修复后路径在入口归一为绝对）"""
        tmp = vt_common.temp_dir('vt_pptrel_')
        try:
            cfg = os.path.join(tmp, 'project_config.json')
            tpl = os.path.join(tmp, 'template.pptx')
            write_config(cfg)
            write_template(tpl)
            rel_out = 'vt_rel_out_tmp'
            p = run_generator(cfg, tpl, rel_out, cwd=tmp)
            assert_deliverable(self, p, os.path.join(tmp, rel_out))
        finally:
            vt_common.rmtree(tmp)


if __name__ == '__main__':
    unittest.main()
