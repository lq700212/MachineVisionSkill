# -*- coding: utf-8 -*-
"""官网产品图获取 回归+注入：横幅合成、型号防串行、诚实降级
（2026-08-30 新功能；注入用例对应"禁止编造图片"约定）"""
import os
import unittest

import vt_common
from PIL import Image

from fetch_product_image import (BANNER_H, BANNER_W, BRAND_ADAPTERS,
                                 anchor_part_shapes, compose_banner,
                                 fetch_camera_image, find_local_image,
                                 find_model_variants, match_model_exact,
                                 match_type_image, migrate_part_images,
                                 replace_part_images, replace_ppt_camera_image,
                                 resolve_selection_parts, sanitize_filename,
                                 scan_part_images, wrap_text)


def _mk_product_img(size=(360, 360)):
    """测试用伪产品图（透明底方块）"""
    img = Image.new('RGBA', size, (0, 0, 0, 0))
    for x in range(size[0] // 4):
        for y in range(size[1] // 4):
            img.putpixel((x * 2, y * 2), (40, 40, 40, 255))
    return img


class TestBannerCompose(unittest.TestCase):
    """横幅合成（离线）"""

    def test_banner_basic(self):
        """回归: 合成输出662x163、文件存在"""
        out = os.path.join(vt_common.temp_dir(), 'm1.png')
        compose_banner('MV-TEST-100GM', '500万像素测试相机，黑白', _mk_product_img(), out)
        self.assertTrue(os.path.exists(out), '横幅未生成')
        with Image.open(out) as im:
            self.assertEqual(im.size, (BANNER_W, BANNER_H),
                             f'横幅尺寸{im.size}≠({BANNER_W},{BANNER_H})')

    def test_banner_overlong_desc_truncated(self):
        """注入: 200字超长描述 → 截断不溢出画布、不抛异常"""
        out = os.path.join(vt_common.temp_dir(), 'm2.png')
        compose_banner('MV-TEST-200GM', '超长' * 100, _mk_product_img(), out)
        with Image.open(out) as im:
            self.assertEqual(im.size, (BANNER_W, BANNER_H))

    def test_wrap_text_lines_within_width(self):
        """回归: 断行后每行渲染宽度不超上限"""
        from PIL import ImageDraw, ImageFont
        img = Image.new('RGB', (800, 200))
        d = ImageDraw.Draw(img)
        f = ImageFont.truetype(r'C:\Windows\Fonts\msyh.ttc', 14)
        lines = wrap_text(d, '很长的中文描述' * 30, f, 350)
        self.assertGreater(len(lines), 1, '长文本应断成多行')
        for ln in lines:
            self.assertLessEqual(d.textlength(ln, font=f), 351,
                                 msg=f'断行溢出: {ln!r}')

    def test_sanitize_filename(self):
        """注入: 型号含非法路径字符 → 替换而非崩溃/目录穿越"""
        self.assertEqual(sanitize_filename('MV-A/B:C*D'), 'MV-A-B-C-D')
        self.assertNotIn('/', sanitize_filename('../../etc'))
        self.assertNotIn('\\', sanitize_filename('a\\b'))


class TestModelMatch(unittest.TestCase):
    """型号精确匹配防串行（对应v1.6.1串行事故：WWK05拿到WWK03参数）"""

    ITEMS = [{'productModel': 'MV-CS050-10GM', 'id': 7483},
             {'productModel': 'MV-CS050-10GM-PRO', 'id': 7324},
             {'productModel': 'MV-CS050-10GM V5', 'id': 9001}]

    def test_exact_match_each_variant(self):
        """回归: 基础版/-PRO/V5 各自精确命中，绝不互相串"""
        for it in self.ITEMS:
            hit = match_model_exact(self.ITEMS, it['productModel'])
            self.assertEqual(hit['id'], it['id'],
                             f"{it['productModel']} 串行到 {hit}")

    def test_match_missing_returns_none(self):
        """注入: 目标不在列表 → None（上层诚实报错），禁止就近匹配"""
        self.assertIsNone(match_model_exact(self.ITEMS, 'MV-CS050-10GM-NPOE'))
        self.assertIsNone(match_model_exact(self.ITEMS, 'MV-CS050-10gm'))  # 大小敏感

    def test_find_variants_prefix_only(self):
        """回归: 变体发现只认同前缀（GM），不混入GC/UM等其它后缀款
        （对应2026-08-30实测: MV-CS060-10GM基础版官网下架，在售仅-PRO/V5）"""
        items = [{'productModel': m} for m in
                 ['MV-CS060-10GM-PRO', 'MV-CS060-10GM V5',
                  'MV-CS060-10GC-PRO', 'MV-CS060-10UM', 'MV-CS060-10GM-PRO-NPOE']]
        self.assertEqual(find_model_variants(items, 'MV-CS060-10GM'),
                         ['MV-CS060-10GM-PRO', 'MV-CS060-10GM V5',
                          'MV-CS060-10GM-PRO-NPOE'])


@unittest.skipUnless(vt_common.online(), '官网不可达，在线用例跳过')
class TestOnlineFetch(unittest.TestCase):
    """在线端到端+注入（真实官网API）"""

    def test_resolve_series(self):
        """回归: 短系列名 → 官网全名唯一命中"""
        from fetch_product_image import resolve_official_series
        name = resolve_official_series(BRAND_ADAPTERS['海康威视'], 'CS系列')
        self.assertEqual(name, 'CS系列工业面阵相机')

    def test_e2e_fetch_to_temp_db(self):
        """回归: 端到端补图——临时库副本+临时输出，出图662x163且写回临时库
        （force=True: 临时库继承的旧图会触发复用短路，必须强制走真实网络链路）"""
        import json
        tmp = vt_common.temp_dir()
        db = vt_common.copy_real_db(os.path.join(tmp, 'db.json'))
        r = fetch_camera_image('海康威视', 'MV-CS016-10GM', db_path=db,
                               out_dir=tmp, force=True)
        self.assertTrue(r['ok'], f'端到端失败: {r["message"]}')
        self.assertTrue(r['db_updated'], 'force全链路应写回临时库')
        with Image.open(r['image_path']) as im:
            self.assertEqual(im.size, (BANNER_W, BANNER_H))
        entry = next(c for c in json.load(
            open(db, encoding='utf-8'))['cameras']
            if c['model'] == 'MV-CS016-10GM')
        self.assertTrue(entry.get('image_path'),
                        '临时库未写回image_path')
        self.assertTrue(entry.get('source_product_id'),
                        'productId未入库（无法追溯）')

    def test_reuse_existing_local_banner(self):
        """回归: 库无条目但本地横幅已在 → 直接复用不重复请求官网（高频场景省时）"""
        r = fetch_camera_image('海康威视', 'MV-CS016-10GM',
                               db_path=os.path.join(vt_common.temp_dir(), 'none.json'))
        self.assertTrue(r['ok'])
        self.assertIn('复用', r['message'])
        self.assertFalse(r['db_updated'])

    def test_inject_model_not_exist(self):
        """注入: 编造型号 → 诚实报错含"转人工"，禁止编造图片"""
        r = fetch_camera_image('海康威视', 'MV-CS999-99XX',
                               db_path=os.path.join(vt_common.temp_dir(), 'none.json'),
                               out_dir=vt_common.temp_dir())
        self.assertFalse(r['ok'], '编造型号竟然出图成功（严重）')
        self.assertIn('转人工', r['message'])

    def test_inject_brand_not_supported(self):
        """注入: 未接入品牌 → 明说"未接入"，禁止假装能取"""
        r = fetch_camera_image('巴斯勒', 'acA1300-30gm',
                               out_dir=vt_common.temp_dir())
        self.assertFalse(r['ok'])
        self.assertIn('未接入', r['message'])

    def test_inject_series_mismatch(self):
        """注入: 系列名错误 → 报命中0个转人工，不猜测"""
        r = fetch_camera_image('海康威视', 'MV-CS016-10GM', series='不存在系列',
                               db_path=os.path.join(vt_common.temp_dir(), 'none.json'),
                               out_dir=vt_common.temp_dir())
        self.assertFalse(r['ok'])
        self.assertIn('转人工', r['message'])

    def test_variant_disambiguation_online(self):
        """注入(在线防串行): -PRO 变体必须拿到与基础版不同的productId，
        绝不串到基础版（对应型号名推断/串行事故）。
        force=True: 临时库继承的旧图会触发复用短路，必须强制现取现验。"""
        import json
        tmp = vt_common.temp_dir()
        db = vt_common.copy_real_db(os.path.join(tmp, 'db.json'))
        r = fetch_camera_image('海康威视', 'MV-CS050-10GM-PRO', db_path=db,
                               out_dir=tmp, force=True)
        self.assertTrue(r['ok'], f"-PRO 获取失败: {r['message']}")
        cams = {c['model']: c for c in
                json.load(open(db, encoding='utf-8'))['cameras']}
        pro_id = cams['MV-CS050-10GM-PRO'].get('source_product_id')
        base_id = cams['MV-CS050-10GM'].get('source_product_id')
        self.assertTrue(pro_id, '-PRO的productId未沉淀')
        self.assertNotEqual(pro_id, base_id,
                            f'-PRO串行到基础版(productId均为{pro_id})')

    def test_missing_base_model_lists_variants(self):
        """回归: 基础版不在售 → 报错信息列出官网同款在售变体
        （2026-08-30实测: MV-CS060-10GM基础版下架，官网仅-PRO/V5在售）"""
        r = fetch_camera_image('海康威视', 'MV-CS060-10GM',
                               db_path=os.path.join(vt_common.temp_dir(), 'none.json'),
                               out_dir=vt_common.temp_dir())
        self.assertFalse(r['ok'], '基础版不在售却出图成功（可能官网已恢复，需更新本用例）')
        self.assertIn('MV-CS060-10GM-PRO', r['message'], '报错未列在售变体')
        self.assertIn('MV-CS060-10GM V5', r['message'])


class TestPptReplace(unittest.TestCase):
    """PPT相机图替换（replace子命令，2026-08-30固化；离线--image直给不联网）"""

    def _mk_pptx(self, path, banner, pages=2):
        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation()
        for _ in range(pages):
            s = prs.slides.add_slide(prs.slide_layouts[6])
            s.shapes.add_picture(banner, Emu(7021859), Emu(638286),
                                 Emu(4736952), Emu(1166349))
        prs.save(path)

    def test_replace_all_pages_original_untouched(self):
        """回归: 全部命中页替换、原PPT字节不动（预览件绝不改原方案）"""
        from pptx import Presentation
        tmp = vt_common.temp_dir()
        old_p = os.path.join(tmp, 'old.png')
        new_p = os.path.join(tmp, 'new.png')
        Image.new('RGB', (662, 170), (200, 210, 220)).save(old_p)
        Image.new('RGB', (BANNER_W, BANNER_H), (240, 120, 120)).save(new_p)
        src = os.path.join(tmp, 'plan.pptx')
        self._mk_pptx(src, old_p)
        src_bytes = open(src, 'rb').read()
        out = os.path.join(tmp, 'out.pptx')
        rep, _ = replace_ppt_camera_image(src, new_p, old_p, out)
        self.assertEqual(len(rep), 2, f'两页应各替换1处: {rep}')
        self.assertTrue(os.path.exists(out), '替换后未产出文件')
        self.assertEqual(open(src, 'rb').read(), src_bytes, '原PPT被改动（严重）')
        blobs = [sh.image.blob for sl in Presentation(out).slides
                 for sh in sl.shapes if sh.shape_type == 13]
        self.assertEqual(blobs, [open(new_p, 'rb').read()] * 2, '新图未全部命中')

    def test_zero_hit_no_output(self):
        """注入: 旧图与PPT内容不符 → 0替换、不留半成品、清单可人工定位"""
        tmp = vt_common.temp_dir()
        old_p = os.path.join(tmp, 'o.png')
        other_p = os.path.join(tmp, 'x.png')
        new_p = os.path.join(tmp, 'n.png')
        Image.new('RGB', (662, 170), (1, 2, 3)).save(old_p)
        Image.new('RGB', (662, 170), (4, 5, 6)).save(other_p)
        Image.new('RGB', (BANNER_W, BANNER_H), (7, 8, 9)).save(new_p)
        src = os.path.join(tmp, 'plan.pptx')
        self._mk_pptx(src, other_p)
        out = os.path.join(tmp, 'out.pptx')
        rep, inv = replace_ppt_camera_image(src, new_p, old_p, out)
        self.assertEqual(rep, [])
        self.assertFalse(os.path.exists(out), '0命中不得留半成品文件')
        self.assertGreaterEqual(len(inv), 2, '应列出图片清单供人工定位')

    def test_cli_replace_end_to_end(self):
        """回归(CLI): 弱模型实际入口——replace --image --old-image 全离线跑通"""
        import subprocess
        import sys
        tmp = vt_common.temp_dir()
        old_p = os.path.join(tmp, 'old.png')
        new_p = os.path.join(tmp, 'new.png')
        Image.new('RGB', (662, 170), (10, 20, 30)).save(old_p)
        Image.new('RGB', (BANNER_W, BANNER_H), (30, 20, 10)).save(new_p)
        src = os.path.join(tmp, 'plan.pptx')
        self._mk_pptx(src, old_p, pages=1)
        out = os.path.join(tmp, 'cli_out.pptx')
        r = subprocess.run(
            [sys.executable,
             os.path.join(vt_common.TOOLS_DIR, 'fetch_product_image.py'),
             'replace', '--pptx', src, '--image', new_p,
             '--old-image', old_p, '--out', out],
            capture_output=True, text=True, timeout=120, cwd=tmp)
        self.assertEqual(r.returncode, 0, f'CLI失败: {r.stdout}{r.stderr}')
        self.assertIn('共1处', r.stdout)
        self.assertIn('原文件未改动', r.stdout)
        self.assertTrue(os.path.exists(out))


class TestReplaceParts(unittest.TestCase):
    """硬件页三件套替换（镜头/光源文字锚定，2026-08-31固化；全程离线+临时库隔离）"""

    @staticmethod
    def _mk_part_pptx(path, cam_png, lens_png, light_png, with_labels=True):
        """构造硬件页布局：右上相机横幅 + (镜头图|镜头文本框) + (光源图|光源文本框)"""
        from pptx import Presentation
        from pptx.util import Emu
        prs = Presentation()
        s = prs.slides.add_slide(prs.slide_layouts[6])
        s.shapes.add_picture(cam_png, Emu(6400000), Emu(500000),
                             Emu(2700000), Emu(660000))
        s.shapes.add_picture(lens_png, Emu(5500000), Emu(1800000),
                             Emu(1200000), Emu(900000))
        s.shapes.add_picture(light_png, Emu(5500000), Emu(3200000),
                             Emu(1200000), Emu(900000))
        if with_labels:
            tb1 = s.shapes.add_textbox(Emu(6800000), Emu(2000000),
                                       Emu(2300000), Emu(500000))
            tb1.text_frame.text = '远心镜头：TEST-LENS-01，0.259x'
            tb2 = s.shapes.add_textbox(Emu(6800000), Emu(3400000),
                                       Emu(2300000), Emu(500000))
            tb2.text_frame.text = 'LED环形光光源：TEST-LIGHT-01'
        prs.save(path)
        return prs

    def setUp(self):
        self.tmp = vt_common.temp_dir()
        self.cam = os.path.join(self.tmp, 'cam.png')
        self.lens_old = os.path.join(self.tmp, 'lens_old.png')
        self.lens_new = os.path.join(self.tmp, 'lens_new.png')
        self.light_old = os.path.join(self.tmp, 'light_old.png')
        self.light_new = os.path.join(self.tmp, 'light_new.png')
        Image.new('RGB', (662, 163), (9, 9, 9)).save(self.cam)
        Image.new('RGB', (400, 300), (1, 2, 3)).save(self.lens_old)
        Image.new('RGB', (500, 400), (11, 22, 33)).save(self.lens_new)
        Image.new('RGB', (300, 300), (4, 5, 6)).save(self.light_old)
        Image.new('RGB', (600, 450), (44, 55, 66)).save(self.light_new)

    def test_anchor_locate(self):
        """回归: 文字锚定命中描述文本框左侧同行的图片；光路标签(无图)不误锚"""
        src = os.path.join(self.tmp, 'plan.pptx')
        self._mk_part_pptx(src, self.cam, self.lens_old, self.light_old)
        from pptx import Presentation
        hits = anchor_part_shapes(Presentation(src),
                                  ('远心镜头', '镜头'))
        self.assertEqual(len(hits), 1, f'镜头锚定应唯一命中: {hits}')
        self.assertEqual(hits[0][0], 1)
        hits_l = anchor_part_shapes(Presentation(src),
                                    ('光源', '环光', '同轴'))
        self.assertEqual(len(hits_l), 1, f'光源锚定应唯一命中: {hits_l}')

    def test_replace_parts_in_place(self):
        """回归: 镜头/光源各自原位替换、contain不变形、原文件字节不动"""
        src = os.path.join(self.tmp, 'plan.pptx')
        self._mk_part_pptx(src, self.cam, self.lens_old, self.light_old)
        src_bytes = open(src, 'rb').read()
        out = os.path.join(self.tmp, 'out.pptx')
        detail, inventory, hits = replace_part_images(
            src, {'lens': self.lens_new, 'light': self.light_new}, out)
        self.assertEqual(hits, {'lens': 1, 'light': 1}, f'各应命中1处: {hits}')
        self.assertTrue(os.path.exists(out))
        self.assertEqual(open(src, 'rb').read(), src_bytes, '原PPT被改动（严重）')
        from pptx import Presentation
        blobs = [sh.image.blob for sl in Presentation(out).slides
                 for sh in sl.shapes if sh.shape_type == 13]
        self.assertIn(open(self.lens_new, 'rb').read(), blobs, '镜头新图未入')
        self.assertIn(open(self.light_new, 'rb').read(), blobs, '光源新图未入')
        self.assertIn(open(self.cam, 'rb').read(), blobs, '相机图不应被动')

    def test_label_without_picture_no_hit(self):
        """注入: 只有描述文本框、无左侧图片 → 0命中不留半成品（不瞎换）"""
        src = os.path.join(self.tmp, 'plan.pptx')
        self._mk_part_pptx(src, self.cam, self.lens_old, self.light_old,
                           with_labels=False)
        out = os.path.join(self.tmp, 'out.pptx')
        detail, inventory, hits = replace_part_images(
            src, {'lens': self.lens_new, 'light': self.light_new}, out)
        self.assertEqual(hits, {'lens': 0, 'light': 0})
        self.assertFalse(os.path.exists(out), '0命中不得留半成品')

    def test_find_local_image_db_first(self):
        """回归: 库image_path显式映射优先；未登记型号返回None不猜"""
        db = os.path.join(self.tmp, 'db.json')
        import json
        json.dump({'cameras': [], 'lenses': [
            {'model': 'TEST-LENS-01',
             'image_path': self.lens_new.replace(os.sep, '/')}]},
            open(db, 'w', encoding='utf-8'))
        self.assertEqual(find_local_image('TEST-LENS-01', db_path=db),
                         self.lens_new)
        self.assertIsNone(find_local_image('TEST-LENS-404', db_path=db),
                          '不存在的型号不得就近猜图')

    def test_type_level_image_mapping(self):
        """回归: 类型图动态扫描——资产文件名与选型type互包含或别名命中
        （用户资产形态：assets/part_images/下中文命名，新增图即自动参与匹配）"""
        tdir = os.path.join(self.tmp, 'timg')
        os.makedirs(tdir, exist_ok=True)
        pics = {}
        for name in ('远心镜头', '定焦镜头', '环光', '同轴光', '条光', '镜头'):
            p = os.path.join(tdir, name + '.png')
            Image.new('RGB', (100, 80), (5, 5, 5)).save(p)
            pics[name] = p
        empty_db = os.path.join(self.tmp, 'none.json')
        # 文件名互包含：type"双远心镜头"含文件名"远心镜头"
        p, how = match_type_image('lens', '双远心镜头', img_dirs=[tdir])
        self.assertEqual(p, pics['远心镜头'], f'远心应命中: {how}')
        p, _ = match_type_image('lens', '定焦镜头', img_dirs=[tdir])
        self.assertEqual(p, pics['定焦镜头'])
        # 别名命中（内置表）：type"LED环形光光源"→环光；"LED条形光光源"→条光
        p, how = match_type_image('light', 'LED环形光光源', img_dirs=[tdir])
        self.assertEqual(p, pics['环光'], f'环光别名应命中: {how}')
        p, how = match_type_image('light', 'LED条形光光源', img_dirs=[tdir])
        self.assertEqual(p, pics['条光'], f'条光别名应命中: {how}')
        p, _ = match_type_image('light', 'LED同轴光光源', img_dirs=[tdir])
        self.assertEqual(p, pics['同轴光'])
        # 陌生type：光源部件无任何关键词沾边 → 不硬配；
        # '广角镜头'命中泛图"镜头.png"是预期行为（泛图兜底），只断言命中泛图而非专属图
        self.assertIsNone(match_type_image('light', '紫外固化灯', img_dirs=[tdir])[0])
        p, _ = match_type_image('lens', '广角镜头', img_dirs=[tdir])
        self.assertEqual(os.path.splitext(os.path.basename(p))[0], '镜头',
                         '广角镜头应落到泛图')
        self.assertIsNone(match_type_image(None, '远心镜头', img_dirs=[tdir])[0],
                          'part为空不得匹配（防相机误配）')
        # 多命中取最具体（stem最长）：泛图"镜头"不得盖过"远心镜头"
        p, how = match_type_image('lens', '双远心镜头', img_dirs=[tdir])
        self.assertEqual(os.path.splitext(os.path.basename(p))[0], '远心镜头',
                          f'应取最具体的类型图: {how}')

    def test_alias_custom_register(self):
        """回归: map --add 等价的别名登记——自定义别名表实时生效，不改代码"""
        tdir = os.path.join(self.tmp, 'timg2')
        os.makedirs(tdir, exist_ok=True)
        tiao = os.path.join(tdir, '条光.png')
        Image.new('RGB', (100, 80), (6, 6, 6)).save(tiao)
        # 无匹配别名：type写法对不上 → None（不硬配）
        self.assertIsNone(match_type_image('light', 'LED条形光光源',
                                           img_dirs=[tdir], aliases={})[0])
        # 登记别名后（等价 map --add 条光 --alias 条形光）立即命中
        p, how = match_type_image('light', 'LED条形光光源', img_dirs=[tdir],
                                  aliases={'条光': ['条形光']})
        self.assertEqual(p, tiao, f'登记别名后应命中: {how}')

    def test_migrate_part_images(self):
        """回归: --migrate 搬家——中文命名类型图迁走，型号横幅留在images/"""
        src = os.path.join(self.tmp, 'm_src')
        dst = os.path.join(self.tmp, 'm_dst')
        os.makedirs(src, exist_ok=True)
        Image.new('RGB', (662, 163), (8, 8, 8)).save(os.path.join(src, 'MV-CS050-10GM.png'))
        Image.new('RGB', (100, 80), (9, 9, 9)).save(os.path.join(src, '远心镜头.png'))
        Image.new('RGB', (100, 80), (10, 10, 10)).save(os.path.join(src, '条光.png'))
        moved = migrate_part_images(src_dir=src, dst_dir=dst)
        self.assertEqual(len(moved), 2, f'应迁移2张类型图: {moved}')
        self.assertTrue(os.path.exists(os.path.join(dst, '远心镜头.png')))
        self.assertTrue(os.path.exists(os.path.join(dst, '条光.png')))
        self.assertFalse(os.path.exists(os.path.join(src, '远心镜头.png')),
                         '类型图迁走后不应留在images/')
        self.assertTrue(os.path.exists(os.path.join(src, 'MV-CS050-10GM.png')),
                        '相机横幅不得被迁移')

    def test_scan_and_find_combined(self):
        """回归: find_local_image 三级链——库>型号图>类型图，扫描目录同时看资产+images兼容"""
        tdir = os.path.join(self.tmp, 'timg3')
        os.makedirs(tdir, exist_ok=True)
        yx = os.path.join(tdir, '远心镜头.png')
        Image.new('RGB', (100, 80), (11, 11, 11)).save(yx)
        empty_db = os.path.join(self.tmp, 'none.json')
        self.assertEqual(find_local_image('TEST-L9', db_path=empty_db, part='lens',
                                          type_str='双远心镜头', img_dirs=[tdir]), yx)
        self.assertIsNone(find_local_image('TEST-L9', db_path=empty_db, part=None,
                                           type_str='双远心镜头', img_dirs=[tdir]))
        stems = scan_part_images([tdir])
        self.assertIn('远心镜头', stems)

    def test_resolve_selection_parts_degrades(self):
        """回归: selection_result三件型号解析；缺图部件None+WARN说明，不阻断"""
        import json
        tmp2 = os.path.join(self.tmp, 'ws')
        os.makedirs(tmp2, exist_ok=True)
        db = os.path.join(self.tmp, 'db.json')
        json.dump({'cameras': [{'model': 'MV-TEST-CAM',
                                'image_path': self.cam.replace(os.sep, '/')}],
                   'lenses': [], 'light_sources': []},
                  open(db, 'w', encoding='utf-8'))
        sel = os.path.join(tmp2, 'selection_result.json')
        json.dump({'hardware_selection': {
            'camera': {'model': 'MV-TEST-CAM'},
            'lens': {'model': 'TEST-LENS-404'},
            'light_source': {'model': 'TEST-LIGHT-404'}}},
            open(sel, 'w', encoding='utf-8'))
        parts, notes = resolve_selection_parts(sel, db_path=db)
        self.assertEqual(parts['camera'], self.cam)
        self.assertIsNone(parts['lens'], '缺图镜头应为None')
        self.assertIsNone(parts['light'])
        self.assertTrue(any('无匹配产品图' in n for n in notes),
                        f'缺图应有WARN说明: {notes}')
        self.assertTrue(any('map --add' in n for n in notes),
                        f'缺图提示应带实时登记指引: {notes}')

    def test_cli_replace_all_parts_end_to_end(self):
        """回归(CLI): 弱模型入口——replace --all_parts 三件套一条命令离线跑通"""
        import json
        import subprocess
        import sys
        tmp = self.tmp
        db = os.path.join(tmp, 'db.json')
        json.dump({'cameras': [{'model': 'MV-TEST-CAM',
                                'image_path': self.cam.replace(os.sep, '/')}],
                   'lenses': [{'model': 'TEST-LENS-01',
                               'image_path': self.lens_new.replace(os.sep, '/')}],
                   'light_sources': [{'model': 'TEST-LIGHT-01',
                                      'image_path': self.light_new.replace(os.sep, '/')}]},
                  open(db, 'w', encoding='utf-8'))
        src = os.path.join(tmp, 'plan.pptx')
        self._mk_part_pptx(src, self.cam, self.lens_old, self.light_old)
        # 真实场景约定：selection_result.json 与 PPT 同目录（主流程产出）
        sel = os.path.join(tmp, 'selection_result.json')
        json.dump({'hardware_selection': {
            'camera': {'model': 'MV-TEST-CAM'},
            'lens': {'model': 'TEST-LENS-01', 'type': '双远心镜头'},
            'light_source': {'model': 'TEST-LIGHT-01', 'type': 'LED环形光光源'}}},
            open(sel, 'w', encoding='utf-8'))
        r = subprocess.run(
            [sys.executable,
             os.path.join(vt_common.TOOLS_DIR, 'fetch_product_image.py'),
             'replace', '--pptx', src, '--all_parts', '--db', db],
            capture_output=True, text=True, timeout=120, cwd=tmp)
        self.assertEqual(r.returncode, 0, f'CLI失败: {r.stdout}{r.stderr}')
        for frag in ('相机1处', '镜头1处', '光源1处', '硬件图预览就绪'):
            self.assertIn(frag, r.stdout, f'汇总缺[{frag}]: {r.stdout}')
        out = os.path.join(tmp, 'plan_硬件图预览.pptx')
        self.assertTrue(os.path.exists(out))


if __name__ == '__main__':
    unittest.main()
